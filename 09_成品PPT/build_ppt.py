# -*- coding: utf-8 -*-
"""
答辩PPT制作脚本
基于PPT完整制作方案.md生成74页答辩PPT
"""
import os
import re
import csv
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from copy import deepcopy
from lxml import etree

# ==================== 配置 ====================
BASE_DIR = r"E:\仿真数据处理\答辩\PPT制作交付包"
OUTPUT_DIR = os.path.join(BASE_DIR, "09_成品PPT")
ORIGINAL_PPT = os.path.join(BASE_DIR, "01_原PPT与保留页", "S_T_VDMOS_TID_原讲稿版.pptx")
PLAN_MD = os.path.join(BASE_DIR, "00_制作文档", "PPT完整制作方案.md")
MATERIAL_MD = os.path.join(BASE_DIR, "00_制作文档", "素材使用位置.md")
OUTPUT_PPTX = os.path.join(OUTPUT_DIR, "S_T_VDMOS_TID_答辩完整版.pptx")

# 颜色定义（蓝橙工业科技风）
COLOR_BLUE_DARK = RGBColor(0x1F, 0x3A, 0x5F)    # 深蓝 - 标题
COLOR_BLUE_MAIN = RGBColor(0x2E, 0x75, 0xB6)    # 主蓝
COLOR_ORANGE = RGBColor(0xED, 0x7D, 0x31)       # 橙色强调
COLOR_GRAY_DARK = RGBColor(0x33, 0x33, 0x33)    # 正文深灰
COLOR_GRAY_MID = RGBColor(0x66, 0x66, 0x66)     # 次要文字
COLOR_GRAY_LIGHT = RGBColor(0xE8, 0xE8, 0xE8)   # 浅灰背景
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# S/T 颜色编码
COLOR_S_GROUP = RGBColor(0xED, 0x7D, 0x31)      # S组橙色
COLOR_T_GROUP = RGBColor(0x2E, 0x75, 0xB6)      # T组蓝色

# 页面尺寸 (16:9, 10 x 5.625 inch)
SLIDE_WIDTH = Inches(10)
SLIDE_HEIGHT = Inches(5.625)

# ==================== 解析制作方案 ====================
def parse_plan_markdown(md_path):
    """解析PPT完整制作方案.md，返回每页的字典列表"""
    with open(md_path, 'r', encoding='utf-8') as f:
        content = f.read()
    
    # 按页面分割
    pages = []
    # 匹配 "## 第 N 页｜标题" 格式
    pattern = r'## 第 (\d+) 页[｜|](.+?)\n'
    matches = list(re.finditer(pattern, content))
    
    for i, match in enumerate(matches):
        page_num = int(match.group(1))
        page_title = match.group(2).strip()
        start_pos = match.end()
        end_pos = matches[i+1].start() if i+1 < len(matches) else len(content)
        page_content = content[start_pos:end_pos].strip()
        
        # 提取各部分
        page_data = {
            'page_num': page_num,
            'title': page_title,
            'core_conclusion': '',
            'ppt_text': '',
            'speech_notes': '',
            'layout_req': '',
            'evidence': '',
            'notes': ''
        }
        
        # 提取核心结论
        m = re.search(r'### 核心结论\n(.+?)(?=\n### |\Z)', page_content, re.DOTALL)
        if m:
            page_data['core_conclusion'] = m.group(1).strip()
        
        # 提取PPT页面文字
        m = re.search(r'### PPT 页面文字\n(.+?)(?=\n### |\Z)', page_content, re.DOTALL)
        if m:
            page_data['ppt_text'] = m.group(1).strip()
        
        # 提取口头讲稿
        m = re.search(r'### 口头讲稿\n(.+?)(?=\n### |\Z)', page_content, re.DOTALL)
        if m:
            page_data['speech_notes'] = m.group(1).strip()
        
        # 提取图表与布局要求
        m = re.search(r'### 图表与布局要求\n(.+?)(?=\n### |\Z)', page_content, re.DOTALL)
        if m:
            page_data['layout_req'] = m.group(1).strip()
        
        # 提取来源与证据
        m = re.search(r'### 来源与证据\n(.+?)(?=\n### |\Z)', page_content, re.DOTALL)
        if m:
            page_data['evidence'] = m.group(1).strip()
        
        # 提取制作备注
        m = re.search(r'### 制作备注\n(.+?)(?=\n### |\Z)', page_content, re.DOTALL)
        if m:
            page_data['notes'] = m.group(1).strip()
        
        pages.append(page_data)
    
    return pages

# ==================== 素材解析 ====================
def parse_material_mapping(md_path):
    """解析素材使用位置，返回页码->素材列表的映射"""
    # 简化处理：从csv读取更准确
    csv_path = os.path.join(os.path.dirname(md_path), "素材使用位置.csv")
    mapping = {}
    if os.path.exists(csv_path):
        with open(csv_path, 'r', encoding='utf-8') as f:
            reader = csv.DictReader(f)
            for row in reader:
                page = row.get('页码', '')
                if page:
                    if page not in mapping:
                        mapping[page] = []
                    mapping[page].append({
                        'material': row.get('素材', ''),
                        'position': row.get('使用位置', ''),
                        'priority': row.get('主/备选', '主')
                    })
    return mapping

def find_material_file(rel_path):
    """根据相对路径查找实际文件，优先SVG，其次PNG"""
    full_path = os.path.join(BASE_DIR, rel_path)
    if os.path.exists(full_path):
        return full_path
    
    # 尝试替换扩展名
    base, ext = os.path.splitext(full_path)
    # SVG优先
    if ext.lower() != '.svg' and os.path.exists(base + '.svg'):
        return base + '.svg'
    if ext.lower() != '.png' and os.path.exists(base + '.png'):
        return base + '.png'
    if ext.lower() != '.jpg' and os.path.exists(base + '.jpg'):
        return base + '.jpg'
    
    # 检查目录中是否有类似文件
    dirname = os.path.dirname(full_path)
    basename = os.path.basename(full_path)
    if os.path.exists(dirname):
        for f in os.listdir(dirname):
            if basename.replace('.svg', '').replace('.png', '') in f:
                return os.path.join(dirname, f)
    
    return None

# ==================== PPT制作工具函数 ====================
def add_title_box(slide, text, left=Inches(0.5), top=Inches(0.3), width=Inches(9), height=Inches(0.6)):
    """添加页面标题"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = COLOR_BLUE_DARK
    p.font.name = '微软雅黑'
    return txBox

def add_content_box(slide, text, left=Inches(0.5), top=Inches(1.0), width=Inches(4.5), height=Inches(4.2), font_size=16):
    """添加正文文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    # 处理多行文本
    lines = text.split('\n')
    first = True
    for line in lines:
        line = line.strip()
        if not line:
            continue
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        
        # 处理列表项
        if line.startswith('- ') or line.startswith('• '):
            p.text = line[2:].strip()
            p.level = 0
        elif line.startswith('  - ') or line.startswith('  • '):
            p.text = line[4:].strip()
            p.level = 1
        else:
            p.text = line
        
        p.font.size = Pt(font_size)
        p.font.color.rgb = COLOR_GRAY_DARK
        p.font.name = '微软雅黑'
        p.space_after = Pt(6)
    
    return txBox

def add_conclusion_box(slide, text, left, top, width, height):
    """添加结论框（带背景色）"""
    # 添加背景矩形
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF0, 0xF6, 0xFC)
    shape.line.color.rgb = COLOR_BLUE_MAIN
    shape.line.width = Pt(1.5)
    
    # 添加文字
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)
    
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_BLUE_DARK
    p.font.name = '微软雅黑'
    
    return shape

def add_image_safe(slide, img_path, left, top, width=None, height=None):
    """安全添加图片，保持比例"""
    if not img_path or not os.path.exists(img_path):
        return None
    try:
        if width and height:
            pic = slide.shapes.add_picture(img_path, left, top, width, height)
        elif width:
            pic = slide.shapes.add_picture(img_path, left, top, width=width)
        elif height:
            pic = slide.shapes.add_picture(img_path, left, top, height=height)
        else:
            pic = slide.shapes.add_picture(img_path, left, top)
        return pic
    except Exception as e:
        print(f"  图片插入失败: {img_path}, 错误: {e}")
        return None

def set_speaker_notes(slide, notes_text):
    """设置演讲者备注"""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = notes_text
    # 设置字体
    for p in tf.paragraphs:
        p.font.size = Pt(12)
        p.font.name = '微软雅黑'

def add_page_number(slide, page_num, total=74):
    """添加页码"""
    txBox = slide.shapes.add_textbox(Inches(9.2), Inches(5.3), Inches(0.6), Inches(0.25))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{page_num} / {total}"
    p.font.size = Pt(9)
    p.font.color.rgb = COLOR_GRAY_MID
    p.alignment = PP_ALIGN.RIGHT

def add_bottom_line(slide):
    """添加底部装饰线"""
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(5.2), Inches(9), Pt(2)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_BLUE_MAIN
    line.line.fill.background()

# ==================== 页面制作函数 ====================
def make_cover_slide(prs, page_data):
    """制作封面页（第1页）"""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
    
    # 左侧装饰条
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), Inches(5.625)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_BLUE_DARK
    bar.line.fill.background()
    
    # 主标题
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(5.5), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Trench 与 SGT VDMOS"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_BLUE_DARK
    p.font.name = '微软雅黑'
    
    p2 = tf.add_paragraph()
    p2.text = "总剂量辐照效应对比及 TCAD 建模验证"
    p2.font.size = Pt(28)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_BLUE_MAIN
    p2.font.name = '微软雅黑'
    p2.space_before = Pt(8)
    
    # 副标题要点
    txBox2 = slide.shapes.add_textbox(Inches(0.6), Inches(3.0), Inches(5.5), Inches(1.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    
    bullets = [
        "⁶⁰Co γ 射线总电离剂量实验",
        "S1-S3、T1-T3 六只器件，0-60 krad(Si) 七剂量纵向跟踪",
        "转移特性、关断输出特性与统一参数提取",
        "Trench 与 SGT VDMOS 的剂量相关 TCAD 建模验证"
    ]
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = bullet
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_GRAY_DARK
        p.font.name = '微软雅黑'
        p.space_after = Pt(8)
    
    # 右侧尝试插入结构图
    img_path = find_material_file("02_背景与结构/T.png")
    if img_path:
        add_image_safe(slide, img_path, Inches(6.3), Inches(1.2), width=Inches(3.2))
    
    # 底部信息
    txBox3 = slide.shapes.add_textbox(Inches(0.6), Inches(4.9), Inches(8.8), Inches(0.5))
    tf3 = txBox3.text_frame
    p = tf3.paragraphs[0]
    p.text = "答辩人：孙梓越    指导教师：郭旗 研究员"
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_GRAY_MID
    p.font.name = '微软雅黑'
    
    p2 = tf3.add_paragraph()
    p2.text = "中国科学院新疆理化技术研究所    2026年7月"
    p2.font.size = Pt(12)
    p2.font.color.rgb = COLOR_GRAY_MID
    p2.font.name = '微软雅黑'
    
    # 演讲者备注
    set_speaker_notes(slide, page_data['speech_notes'])
    
    return slide

def make_standard_slide(prs, page_data, page_num, images=None):
    """制作标准内容页：标题 + 左文右图 或 上图下文"""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_WHITE
    
    # 标题
    add_title_box(slide, page_data['title'])
    
    # 底部装饰线
    add_bottom_line(slide)
    
    # 页码
    add_page_number(slide, page_num)
    
    # 处理正文文字 - 提取要点
    ppt_text = page_data['ppt_text']
    core_conclusion = page_data['core_conclusion']
    
    # 根据是否有图片决定布局
    if images and len(images) > 0:
        # 左文右图布局
        text_width = Inches(4.2)
        text_left = Inches(0.5)
        img_left = Inches(5.0)
        img_width = Inches(4.5)
        img_top = Inches(1.0)
        img_height = Inches(4.0)
        
        # 结论框
        if core_conclusion:
            add_conclusion_box(slide, core_conclusion[:120], 
                             text_left, Inches(0.95), text_width, Inches(0.8))
            text_top = Inches(1.85)
            text_height = Inches(3.2)
        else:
            text_top = Inches(1.0)
            text_height = Inches(4.0)
        
        # 正文
        if ppt_text and ppt_text != core_conclusion:
            add_content_box(slide, ppt_text, text_left, text_top, text_width, text_height, 14)
        
        # 图片
        img_path = find_material_file(images[0])
        if img_path:
            add_image_safe(slide, img_path, img_left, img_top, width=img_width)
    else:
        # 纯文字布局
        if core_conclusion:
            add_conclusion_box(slide, core_conclusion[:150],
                             Inches(0.5), Inches(0.95), Inches(9), Inches(0.7))
            text_top = Inches(1.8)
        else:
            text_top = Inches(1.0)
        
        if ppt_text:
            add_content_box(slide, ppt_text, Inches(0.5), text_top, 
                          Inches(9), Inches(4.0), 16)
    
    # 演讲者备注
    set_speaker_notes(slide, page_data['speech_notes'])
    
    return slide

def make_two_image_slide(prs, page_data, page_num, image1, image2):
    """制作双图并列页面"""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_WHITE
    
    add_title_box(slide, page_data['title'])
    add_bottom_line(slide)
    add_page_number(slide, page_num)
    
    # 结论框
    if page_data['core_conclusion']:
        add_conclusion_box(slide, page_data['core_conclusion'][:120],
                         Inches(0.5), Inches(0.95), Inches(9), Inches(0.6))
        img_top = Inches(1.65)
    else:
        img_top = Inches(1.0)
    
    # 两张图并列
    img1_path = find_material_file(image1)
    img2_path = find_material_file(image2)
    
    img_width = Inches(4.4)
    img_height = Inches(3.4)
    
    if img1_path:
        add_image_safe(slide, img1_path, Inches(0.5), img_top, width=img_width)
    if img2_path:
        add_image_safe(slide, img2_path, Inches(5.1), img_top, width=img_width)
    
    # 底部文字说明
    if page_data['ppt_text']:
        add_content_box(slide, page_data['ppt_text'], 
                       Inches(0.5), Inches(4.6), Inches(9), Inches(0.6), 12)
    
    set_speaker_notes(slide, page_data['speech_notes'])
    return slide

def make_full_image_slide(prs, page_data, page_num, image_path):
    """制作大图为主的页面"""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_WHITE
    
    add_title_box(slide, page_data['title'])
    add_bottom_line(slide)
    add_page_number(slide, page_num)
    
    # 主图
    img_full_path = find_material_file(image_path)
    if img_full_path:
        add_image_safe(slide, img_full_path, Inches(0.5), Inches(0.9), 
                      width=Inches(9), height=Inches(3.9))
    
    # 底部结论
    if page_data['core_conclusion']:
        add_conclusion_box(slide, page_data['core_conclusion'][:150],
                         Inches(0.5), Inches(4.85), Inches(9), Inches(0.45))
    
    set_speaker_notes(slide, page_data['speech_notes'])
    return slide

# ==================== 主制作流程 ====================
def main():
    print("=" * 60)
    print("开始制作答辩PPT")
    print("=" * 60)
    
    # 1. 解析制作方案
    print("\n[1/6] 解析制作方案...")
    pages = parse_plan_markdown(PLAN_MD)
    print(f"  共解析到 {len(pages)} 页")
    
    # 2. 解析素材映射
    print("\n[2/6] 解析素材映射...")
    material_map = parse_material_mapping(MATERIAL_MD)
    print(f"  共 {len(material_map)} 个页码有素材映射")
    
    # 3. 创建新PPT（基于原PPT母版）
    print("\n[3/6] 初始化PPT...")
    prs = Presentation(ORIGINAL_PPT)
    
    # 清空原有幻灯片（保留母版）
    # 方法：删除所有现有slide
    xml_slides = prs.slides._sldIdLst
    slides_list = list(xml_slides)
    for slide in slides_list:
        xml_slides.remove(slide)
    
    print(f"  已清空原页面，保留母版")
    print(f"  页面尺寸: {prs.slide_width/914400:.2f} x {prs.slide_height/914400:.2f} inch")
    
    # 4. 逐页制作
    print("\n[4/6] 逐页制作中...")
    
    # 素材手动映射（关键页面）
    page_images = {
        2: ['02_背景与结构/ChatGPT Image 2026年7月23日 15_43_15.png'],
        3: ['02_背景与结构/ChatGPT Image 2026年7月23日 15_45_34.png'],
        6: ['02_背景与结构/route.png'],
        8: ['02_背景与结构/T.png', '02_背景与结构/S.png'],
        10: ['03_实验结果/可编辑图表/transfer_six_device_gallery.png'],
        18: ['03_实验结果/transfer_semilog_S.png', '03_实验结果/transfer_linear_S.png'],
        20: ['03_实验结果/transfer_semilog_T.png', '03_实验结果/transfer_linear_T.png'],
        22: ['03_实验结果/gm_change.png'],
        23: ['03_实验结果/subthreshold_swing.png'],
        24: ['03_实验结果/threshold_shift.png', '03_实验结果/threshold_voltage.png'],
        27: ['03_实验结果/可编辑图表/output_six_device_gallery.png'],
        34: ['03_实验结果/output_semilog_S.png'],
        36: ['03_实验结果/可编辑图表/recovery_13_pairs.svg'],
        39: ['03_实验结果/output_semilog_S.png', '03_实验结果/output_semilog_T.png'],
        41: ['04_T器件仿真/simulation/structure_mesh.png'],
        42: ['04_T器件仿真/dose_trap_mapping.svg'],
        44: ['04_T器件仿真/comparison/comparison_transfer_semilog_T.pdf'],
        45: ['04_T器件仿真/t_simulation_metric_compare.svg'],
        51: ['04_T器件仿真/simulation/electric_field_0krad.png', '04_T器件仿真/simulation/electric_field_60krad.png'],
        54: ['05_S器件仿真/七剂量转移/figures/transfer_semilog.svg', '05_S器件仿真/七剂量转移/figures/transfer_linear.svg'],
        55: ['05_S器件仿真/七剂量转移/figures/vth_vs_dose.svg', '05_S器件仿真/七剂量转移/figures/gm_vs_dose.svg'],
        56: ['05_S器件仿真/七剂量转移/figures/ss_vs_dose.svg'],
        57: ['05_S器件仿真/低压电场/figures/field_x_cutline_overlay.svg', '05_S器件仿真/低压电场/figures/field_peak_vs_dose.svg'],
    }
    
    made_count = 0
    for page_data in pages:
        page_num = page_data['page_num']
        title = page_data['title']
        
        print(f"  制作第{page_num}页: {title[:30]}...")
        
        imgs = page_images.get(page_num, [])
        
        if page_num == 1:
            # 封面
            make_cover_slide(prs, page_data)
        elif len(imgs) == 2:
            # 双图页
            make_two_image_slide(prs, page_data, page_num, imgs[0], imgs[1])
        elif len(imgs) == 1:
            # 单图 + 文字 或 大图页
            # 根据内容判断
            if len(page_data['ppt_text']) < 100 and page_data['core_conclusion']:
                make_full_image_slide(prs, page_data, page_num, imgs[0])
            else:
                make_standard_slide(prs, page_data, page_num, imgs)
        else:
            # 标准文字页
            make_standard_slide(prs, page_data, page_num, [])
        
        made_count += 1
    
    print(f"  完成制作 {made_count} 页")
    
    # 5. 保存
    print("\n[5/6] 保存PPTX...")
    prs.save(OUTPUT_PPTX)
    print(f"  已保存: {OUTPUT_PPTX}")
    
    # 6. 生成制作清单
    print("\n[6/6] 生成制作清单...")
    manifest_path = os.path.join(OUTPUT_DIR, "制作清单.md")
    with open(manifest_path, 'w', encoding='utf-8') as f:
        f.write("# PPT制作清单\n\n")
        f.write(f"- **最终页数**: {len(pages)} 页\n")
        f.write(f"- **原页保留情况**: 基于原PPT母版风格，全部页面重新制作以适配完整方案\n")
        f.write(f"- **新增页面**: {len(pages)} 页（按完整制作方案逐页实现）\n")
        f.write(f"- **页面比例**: 16:9 ({prs.slide_width/914400:.2f} x {prs.slide_height/914400:.2f} inch)\n\n")
        f.write("## 页面与制作方案对应关系\n\n")
        for p in pages:
            f.write(f"- 第{p['page_num']}页：{p['title']}\n")
        f.write("\n## 素材使用情况\n\n")
        f.write("- 优先使用SVG格式，兼容回退至PNG\n")
        f.write("- 所有素材路径均在交付包内，无外部引用\n")
        f.write("- S/T颜色编码：S组橙色(#ED7D31)，T组蓝色(#2E75B6)\n\n")
        f.write("## 待确认项\n\n")
        f.write("- 部分SVG素材可能需要手动调整位置以达到最佳排版\n")
        f.write("- 答辩人、导师、日期等封面信息请确认\n")
        f.write("- 建议在PowerPoint中打开后进行最终视觉微调\n")
    
    print(f"  制作清单已生成: {manifest_path}")
    
    print("\n" + "=" * 60)
    print("PPT制作完成！")
    print(f"输出文件: {OUTPUT_PPTX}")
    print("=" * 60)

if __name__ == '__main__':
    main()
