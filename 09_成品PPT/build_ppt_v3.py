# -*- coding: utf-8 -*-
"""
答辩PPT制作脚本 v3
创建全新PPT，避免原PPT清空导致的重复名称问题
"""
import os
import sys
import re
import csv

if sys.platform == 'win32':
    import io
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
    sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ==================== 配置 ====================
BASE_DIR = r"E:\仿真数据处理\答辩\PPT制作交付包"
OUTPUT_DIR = os.path.join(BASE_DIR, "09_成品PPT")
PLAN_MD = os.path.join(BASE_DIR, "00_制作文档", "PPT完整制作方案.md")
OUTPUT_PPTX = os.path.join(OUTPUT_DIR, "S_T_VDMOS_TID_答辩完整版.pptx")

# 颜色
COLOR_BLUE_DARK = RGBColor(0x1F, 0x3A, 0x5F)
COLOR_BLUE_MAIN = RGBColor(0x2E, 0x75, 0xB6)
COLOR_ORANGE = RGBColor(0xED, 0x7D, 0x31)
COLOR_GRAY_DARK = RGBColor(0x33, 0x33, 0x33)
COLOR_GRAY_MID = RGBColor(0x66, 0x66, 0x66)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

TOTAL_PAGES = 74

# ==================== 解析方案 ====================
def parse_plan_markdown(md_path):
    with open(md_path, 'r', encoding='utf-8') as f:
        content = f.read()
    
    pages = []
    pattern = r'## 第 (\d+) 页[｜|](.+?)\n'
    matches = list(re.finditer(pattern, content))
    
    for i, match in enumerate(matches):
        page_num = int(match.group(1))
        page_title = match.group(2).strip()
        start_pos = match.end()
        end_pos = matches[i+1].start() if i+1 < len(matches) else len(content)
        page_content = content[start_pos:end_pos].strip()
        
        page_data = {
            'page_num': page_num,
            'title': page_title,
            'core_conclusion': '',
            'ppt_text': '',
            'speech_notes': '',
        }
        
        m = re.search(r'### 核心结论\n(.+?)(?=\n### |\Z)', page_content, re.DOTALL)
        if m:
            page_data['core_conclusion'] = m.group(1).strip()
        
        m = re.search(r'### PPT 页面文字\n(.+?)(?=\n### |\Z)', page_content, re.DOTALL)
        if m:
            page_data['ppt_text'] = m.group(1).strip()
        
        m = re.search(r'### 口头讲稿\n(.+?)(?=\n### |\Z)', page_content, re.DOTALL)
        if m:
            page_data['speech_notes'] = m.group(1).strip()
        
        pages.append(page_data)
    
    return pages

# ==================== 素材映射 ====================
PAGE_IMAGES = {
    2: ['02_背景与结构/ChatGPT Image 2026年7月23日 15_43_15.png'],
    3: ['02_背景与结构/ChatGPT Image 2026年7月23日 15_45_34.png'],
    6: ['02_背景与结构/route.png'],
    7: ['02_背景与结构/site_irradiator.jpg'],
    8: ['02_背景与结构/T.png', '02_背景与结构/S.png'],
    10: ['03_实验结果/可编辑图表/transfer_six_device_gallery.png'],
    11: ['03_实验结果/可编辑图表/transfer_S1.png'],
    12: ['03_实验结果/可编辑图表/transfer_S2.png'],
    13: ['03_实验结果/可编辑图表/transfer_S3.png'],
    14: ['03_实验结果/可编辑图表/transfer_T1.png'],
    15: ['03_实验结果/可编辑图表/transfer_T2.png'],
    16: ['03_实验结果/可编辑图表/transfer_T3.png'],
    17: ['03_实验结果/可编辑图表/transfer_metrics_overview.png'],
    18: ['03_实验结果/transfer_semilog_S.png', '03_实验结果/transfer_linear_S.png'],
    20: ['03_实验结果/transfer_semilog_T.png', '03_实验结果/transfer_linear_T.png'],
    22: ['03_实验结果/gm_change.png'],
    23: ['03_实验结果/subthreshold_swing.png'],
    24: ['03_实验结果/threshold_shift.png', '03_实验结果/threshold_voltage.png'],
    27: ['03_实验结果/可编辑图表/output_six_device_gallery.png'],
    28: ['03_实验结果/可编辑图表/output_S1.png'],
    29: ['03_实验结果/可编辑图表/output_S2.png'],
    30: ['03_实验结果/可编辑图表/output_S3.png'],
    31: ['03_实验结果/可编辑图表/output_T1.png'],
    32: ['03_实验结果/可编辑图表/output_T2.png'],
    33: ['03_实验结果/可编辑图表/output_T3.png'],
    35: ['03_实验结果/可编辑图表/output_metrics_overview.png'],
    34: ['03_实验结果/output_semilog_S.png'],
    36: ['03_实验结果/可编辑图表/recovery_13_pairs.png'],
    37: ['03_实验结果/S_T_VDMOS_TID_完整成果展示版_讲稿版_27.png'],
    39: ['03_实验结果/output_semilog_S.png', '03_实验结果/output_semilog_T.png'],
    41: ['04_T器件仿真/simulation/structure_mesh.png'],
    42: ['03_实验结果/可编辑图表/dose_trap_mapping.png'],
    44: ['04_T器件仿真/comparison/comparison_transfer_semilog_T.png'],
    45: ['03_实验结果/可编辑图表/t_simulation_metric_compare.png'],
    47: ['04_T器件仿真/comparison/comparison_output_semilog_T.png'],
    48: ['04_T器件仿真/comparison/comparison_id30_T.png'],
    49: ['04_T器件仿真/comparison/comparison_high_voltage_context_T.png'],
    51: ['04_T器件仿真/simulation/electric_field_0krad.png', '04_T器件仿真/simulation/electric_field_60krad.png'],
    54: ['05_S器件仿真/七剂量转移/figures/transfer_semilog.png', '05_S器件仿真/七剂量转移/figures/transfer_linear.png'],
    55: ['05_S器件仿真/七剂量转移/figures/vth_vs_dose.png', '05_S器件仿真/七剂量转移/figures/gm_vs_dose.png'],
    56: ['05_S器件仿真/七剂量转移/figures/ss_vs_dose.png'],
    57: ['05_S器件仿真/低压电场/figures/field_x_cutline_overlay.png', '05_S器件仿真/低压电场/figures/field_peak_vs_dose.png'],
    59: ['05_S器件仿真/七剂量转移/figures/svisual_transfer_semilog_0to60K_clean.png'],
    68: ['06_AI辅助流程/AI操作中图像1.png', '06_AI辅助流程/AI操作中图2.png', '06_AI辅助流程/虚拟机工程.png'],
}

def find_image(rel_path):
    if not rel_path:
        return None
    full_path = os.path.join(BASE_DIR, rel_path)
    return full_path if os.path.exists(full_path) else None

# ==================== 辅助函数 ====================
def add_title(slide, text):
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLOR_BLUE_DARK
    p.font.name = '微软雅黑'

def add_bottom_line(slide):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(5.15), Inches(9), Pt(1.5))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_BLUE_MAIN
    line.line.fill.background()

def add_page_num(slide, num):
    txBox = slide.shapes.add_textbox(Inches(9.0), Inches(5.25), Inches(0.7), Inches(0.25))
    p = txBox.text_frame.paragraphs[0]
    p.text = f"{num} / {TOTAL_PAGES}"
    p.font.size = Pt(9)
    p.font.color.rgb = COLOR_GRAY_MID
    p.alignment = PP_ALIGN.RIGHT

def add_conclusion(slide, text, left, top, w, h, fs=12):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xEE, 0xF4, 0xFA)
    shape.line.color.rgb = COLOR_BLUE_MAIN
    shape.line.width = Pt(1.2)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    lines = [l.strip() for l in text.split('\n') if l.strip()]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        clean = line.lstrip('- ').lstrip('• ').strip()
        p.text = clean
        p.font.size = Pt(fs)
        p.font.bold = True
        p.font.color.rgb = COLOR_BLUE_DARK
        p.font.name = '微软雅黑'
        p.space_after = Pt(2)

def add_bullets(slide, text, left, top, w, h, fs=14):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for line in text.split('\n'):
        s = line.strip()
        if not s:
            continue
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        content = s.lstrip('- ').lstrip('• ').strip()
        if s.startswith('  - ') or s.startswith('  • '):
            content = s[4:].strip()
            p.level = 1
        p.text = content
        p.font.size = Pt(fs - (1 if p.level > 0 else 0))
        p.font.color.rgb = COLOR_GRAY_DARK
        p.font.name = '微软雅黑'
        p.space_after = Pt(5)

def add_pic(slide, rel, left, top, width=None, height=None):
    path = find_image(rel)
    if not path:
        return None
    try:
        if width and height:
            return slide.shapes.add_picture(path, left, top, width, height)
        elif width:
            return slide.shapes.add_picture(path, left, top, width=width)
        elif height:
            return slide.shapes.add_picture(path, left, top, height=height)
        else:
            return slide.shapes.add_picture(path, left, top)
    except Exception as e:
        print(f"    [WARN] {rel}: {e}")
        return None

def set_notes(slide, text):
    if not text:
        return
    tf = slide.notes_slide.notes_text_frame
    tf.text = text
    for p in tf.paragraphs:
        p.font.size = Pt(11)
        p.font.name = '微软雅黑'

# ==================== 页面模板 ====================
def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局

def make_cover(prs, pd):
    s = new_slide(prs)
    # 背景白
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = COLOR_WHITE
    
    # 左装饰条
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.12), Inches(5.625))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_BLUE_DARK
    bar.line.fill.background()
    
    # 橙线
    bar2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.9), Inches(2.5), Pt(3))
    bar2.fill.solid()
    bar2.fill.fore_color.rgb = COLOR_ORANGE
    bar2.line.fill.background()
    
    # 标题
    tb = s.shapes.add_textbox(Inches(0.6), Inches(1.05), Inches(5.8), Inches(1.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Trench 与 SGT VDMOS"
    p.font.size = Pt(34); p.font.bold = True; p.font.color.rgb = COLOR_BLUE_DARK; p.font.name = '微软雅黑'
    p2 = tf.add_paragraph()
    p2.text = "总剂量辐照效应对比及 TCAD 建模验证"
    p2.font.size = Pt(26); p2.font.bold = True; p2.font.color.rgb = COLOR_BLUE_MAIN; p2.font.name = '微软雅黑'
    p2.space_before = Pt(6)
    
    # 要点
    tb2 = s.shapes.add_textbox(Inches(0.6), Inches(2.9), Inches(5.5), Inches(1.8))
    tf2 = tb2.text_frame; tf2.word_wrap = True
    bullets = ["⁶⁰Co γ 射线总电离剂量实验", "S1–S3、T1–T3 六只器件，0–60 krad(Si) 七剂量纵向跟踪",
               "转移特性、关断输出特性与统一参数提取", "剂量相关 TCAD 建模与实验—仿真边界验证"]
    for i, b in enumerate(bullets):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = "▸  " + b
        p.font.size = Pt(15); p.font.color.rgb = COLOR_GRAY_DARK; p.font.name = '微软雅黑'
        p.space_after = Pt(8)
    
    # 右图
    add_pic(s, '02_背景与结构/T.png', Inches(6.3), Inches(0.9), width=Inches(3.3))
    
    # 底部
    tb3 = s.shapes.add_textbox(Inches(0.6), Inches(4.8), Inches(8.8), Inches(0.7))
    tf3 = tb3.text_frame
    p = tf3.paragraphs[0]
    p.text = "答辩人：孙梓越"
    p.font.size = Pt(14); p.font.color.rgb = COLOR_GRAY_DARK; p.font.name = '微软雅黑'
    p2 = tf3.add_paragraph()
    p2.text = "指导教师：郭旗 研究员    中国科学院新疆理化技术研究所"
    p2.font.size = Pt(12); p2.font.color.rgb = COLOR_GRAY_MID; p2.font.name = '微软雅黑'; p2.space_before = Pt(4)
    p3 = tf3.add_paragraph()
    p3.text = "2026 年 7 月"
    p3.font.size = Pt(11); p3.font.color.rgb = COLOR_GRAY_MID; p3.font.name = '微软雅黑'; p3.space_before = Pt(4)
    
    set_notes(s, pd['speech_notes'])
    return s

def make_std_left_text_right_img(prs, pd, pn, img):
    s = new_slide(prs)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = COLOR_WHITE
    add_title(s, pd['title']); add_bottom_line(s); add_page_num(s, pn)
    
    tl = Inches(0.5); tw = Inches(4.0)
    il = Inches(4.8); iw = Inches(4.7); it = Inches(0.95)
    
    if pd['core_conclusion']:
        add_conclusion(s, pd['core_conclusion'], tl, Inches(0.95), tw, Inches(0.85), 12)
        tt = Inches(1.9); th = Inches(3.1)
    else:
        tt = Inches(0.95); th = Inches(4.0)
    
    if pd['ppt_text']:
        add_bullets(s, pd['ppt_text'], tl, tt, tw, th, 14)
    if img:
        add_pic(s, img, il, it, width=iw)
    set_notes(s, pd['speech_notes'])
    return s

def make_two_img(prs, pd, pn, i1, i2):
    s = new_slide(prs)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = COLOR_WHITE
    add_title(s, pd['title']); add_bottom_line(s); add_page_num(s, pn)
    
    if pd['core_conclusion']:
        add_conclusion(s, pd['core_conclusion'], Inches(0.5), Inches(0.95), Inches(9), Inches(0.55), 12)
        it = Inches(1.6)
    else:
        it = Inches(0.95)
    
    iw = Inches(4.4)
    add_pic(s, i1, Inches(0.5), it, width=iw)
    add_pic(s, i2, Inches(5.1), it, width=iw)
    
    if pd['ppt_text']:
        add_bullets(s, pd['ppt_text'], Inches(0.5), Inches(4.5), Inches(9), Inches(0.6), 11)
    set_notes(s, pd['speech_notes'])
    return s

def make_big_img(prs, pd, pn, img):
    s = new_slide(prs)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = COLOR_WHITE
    add_title(s, pd['title']); add_bottom_line(s); add_page_num(s, pn)
    
    add_pic(s, img, Inches(0.5), Inches(0.9), width=Inches(9), height=Inches(3.8))
    
    if pd['core_conclusion']:
        add_conclusion(s, pd['core_conclusion'], Inches(0.5), Inches(4.75), Inches(9), Inches(0.45), 12)
    set_notes(s, pd['speech_notes'])
    return s

def make_text_only(prs, pd, pn):
    s = new_slide(prs)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = COLOR_WHITE
    add_title(s, pd['title']); add_bottom_line(s); add_page_num(s, pn)
    
    if pd['core_conclusion']:
        add_conclusion(s, pd['core_conclusion'], Inches(0.5), Inches(0.95), Inches(9), Inches(0.7), 14)
        tt = Inches(1.75)
    else:
        tt = Inches(0.95)
    
    if pd['ppt_text']:
        add_bullets(s, pd['ppt_text'], Inches(0.5), tt, Inches(9), Inches(4.0), 16)
    set_notes(s, pd['speech_notes'])
    return s

def make_three_img(prs, pd, pn, imgs):
    s = new_slide(prs)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = COLOR_WHITE
    add_title(s, pd['title']); add_bottom_line(s); add_page_num(s, pn)
    
    if pd['core_conclusion']:
        add_conclusion(s, pd['core_conclusion'], Inches(0.5), Inches(0.95), Inches(9), Inches(0.55), 12)
        it = Inches(1.6)
    else:
        it = Inches(0.95)
    
    iw = Inches(2.9)
    for i, img in enumerate(imgs[:3]):
        left = Inches(0.5 + i * 3.05)
        add_pic(s, img, left, it, width=iw)
    set_notes(s, pd['speech_notes'])
    return s

def make_thanks(prs, pd):
    s = new_slide(prs)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = COLOR_WHITE
    
    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.5), Inches(10), Pt(2))
    line.fill.solid(); line.fill.fore_color.rgb = COLOR_BLUE_MAIN; line.line.fill.background()
    
    tb = s.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(1.0))
    p = tb.text_frame.paragraphs[0]
    p.text = "感谢各位老师聆听"
    p.font.size = Pt(40); p.font.bold = True; p.font.color.rgb = COLOR_BLUE_DARK
    p.font.name = '微软雅黑'; p.alignment = PP_ALIGN.CENTER
    
    tb2 = s.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(9), Inches(0.6))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = "敬请批评指正"
    p2.font.size = Pt(20); p2.font.color.rgb = COLOR_GRAY_MID
    p2.font.name = '微软雅黑'; p2.alignment = PP_ALIGN.CENTER
    
    set_notes(s, pd['speech_notes'])
    return s

# ==================== 主流程 ====================
def main():
    print("=" * 60)
    print("答辩PPT制作 v3 - 全新构建")
    print("=" * 60)
    
    print("\n[1/4] 解析方案...")
    pages = parse_plan_markdown(PLAN_MD)
    print(f"  {len(pages)} 页")
    
    print("\n[2/4] 创建全新PPT...")
    prs = Presentation()
    # 设置16:9
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    print(f"  尺寸: {prs.slide_width/914400:.2f} x {prs.slide_height/914400:.2f} inch")
    
    print("\n[3/4] 制作页面...")
    ok = 0
    for pd in pages:
        pn = pd['page_num']
        imgs = PAGE_IMAGES.get(pn, [])
        t = pd['title'][:24]
        print(f"  P{pn:02d}: {t}...", end=' ')
        
        try:
            if pn == 1:
                make_cover(prs, pd)
            elif pn == 74:
                make_thanks(prs, pd)
            elif len(imgs) == 3:
                make_three_img(prs, pd, pn, imgs)
            elif len(imgs) == 2:
                make_two_img(prs, pd, pn, imgs[0], imgs[1])
            elif len(imgs) == 1:
                txt_len = len(pd['ppt_text']) + len(pd['core_conclusion'])
                if txt_len > 200:
                    make_std_left_text_right_img(prs, pd, pn, imgs[0])
                else:
                    make_big_img(prs, pd, pn, imgs[0])
            else:
                make_text_only(prs, pd, pn)
            ok += 1
            print("OK")
        except Exception as e:
            print(f"ERR: {e}")
            make_text_only(prs, pd, pn)
            ok += 1
    
    print(f"\n  完成: {ok}/{len(pages)} 页")
    
    print("\n[4/4] 保存...")
    prs.save(OUTPUT_PPTX)
    sz = os.path.getsize(OUTPUT_PPTX) / 1024
    print(f"  已保存: {OUTPUT_PPTX}")
    print(f"  大小: {sz:.0f} KB")
    
    # 验证
    prs2 = Presentation(OUTPUT_PPTX)
    nc = 0
    for sl in prs2.slides:
        if sl.has_notes_slide and sl.notes_slide.notes_text_frame.text.strip():
            nc += 1
    print(f"  验证: {len(prs2.slides)} 页, {nc} 页有备注")
    
    print("\n" + "=" * 60)
    print("制作完成！")
    print("=" * 60)

if __name__ == '__main__':
    main()
