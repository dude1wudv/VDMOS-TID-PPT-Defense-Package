# -*- coding: utf-8 -*-
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt, Inches

p=Path(r"E:\仿真数据处理\答辩\PPT制作交付包\09_成品PPT\S_T_VDMOS_TID_答辩完整版.pptx")
prs=Presentation(p)

replacements={
 "S/T 组 VDMOS 总电离剂量实验与 T 器件 TCAD 验证":"Trench 与 SGT VDMOS 总剂量辐照效应对比及 TCAD 建模验证",
 "0–60 krad(Si) 完整成果展示版":"0–60 krad(Si) 全面答辩完整版",
 "实验事实、数据处理、T 仿真与证据边界。":"实验对比、统一提参、T/S TCAD 建模与证据边界。",
 "S/T VDMOS TID 实验与 T-only TCAD 验证":"Trench / SGT VDMOS TID 实验与 TCAD 验证"
}
for si,s in enumerate(prs.slides,1):
    for sh in s.shapes:
        if not getattr(sh,"has_text_frame",False):
            continue
        for old,new in replacements.items():
            if old in sh.text:
                for pgh in sh.text_frame.paragraphs:
                    for run in pgh.runs:
                        if old in run.text: run.text=run.text.replace(old,new)
        # 复用原页标题：宽标题栏且位于顶部。
        if sh.top < Inches(.95) and sh.width > Inches(5) and sh.text.strip():
            for pgh in sh.text_frame.paragraphs:
                for run in pgh.runs:
                    if run.text.strip() and (run.font.size is None or run.font.size.pt < 22):
                        run.font.size=Pt(22)
        # 页脚、状态与图注最低 10 pt；同步增加过低文本框高度。
        had_small=False
        for pgh in sh.text_frame.paragraphs:
            for run in pgh.runs:
                if run.text.strip() and run.font.size and run.font.size.pt < 10:
                    run.font.size=Pt(10); had_small=True
        if had_small and sh.height < Inches(.18):
            sh.height=Inches(.18)

# 封面长标题保留 24 pt，并扩大为四行安全区；同步下移副标题。
cover=prs.slides[0]
for sh in cover.shapes:
    if not getattr(sh,"has_text_frame",False): continue
    if sh.name=="Text 0":
        sh.height=Inches(1.9)
        for pgh in sh.text_frame.paragraphs:
            for run in pgh.runs:
                if run.text.strip(): run.font.size=Pt(24)
    elif sh.name=="Text 1": sh.top=Inches(3.02)
    elif sh.name=="Text 2": sh.top=Inches(3.58)

# 复用原页的眉题、页码和状态徽标属于导航元素，维持 8 pt；科学边界文字保持 10 pt 并扩容。
reused={11,12,13,14,15,16,17,27,28,29,30,31,32,33,35,36,37,41,42,44,45,46,47,48,50,51,67}
for si in reused:
    s=prs.slides[si-1]
    for sh in s.shapes:
        if not getattr(sh,"has_text_frame",False) or not sh.text.strip(): continue
        if sh.name in {"Text 1","Text 2"}:
            for pgh in sh.text_frame.paragraphs:
                for run in pgh.runs:
                    if run.text.strip(): run.font.size=Pt(8)
        elif sh.name=="Text 7":
            sh.top=min(sh.top, Inches(5.14))
            sh.height=Inches(.25)
            for pgh in sh.text_frame.paragraphs:
                for run in pgh.runs:
                    if run.text.strip(): run.font.size=Pt(8)
        elif sh.name=="Text 8":
            sh.top=min(sh.top, Inches(5.10))
            sh.height=Inches(.46)

prs.save(p)
print('fixed',p)