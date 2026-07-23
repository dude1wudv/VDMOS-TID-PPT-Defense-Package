# -*- coding: utf-8 -*-
from __future__ import annotations

import csv
import hashlib
import json
import os
import re
from pathlib import Path
from typing import Iterable

import fitz
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(r"E:\仿真数据处理\答辩\PPT制作交付包")
OUT = ROOT / "09_成品PPT"
WORK = OUT / "制作工程"
CACHE = WORK / "asset_cache"
PPTX = OUT / "S_T_VDMOS_TID_答辩完整版_底稿.pptx"

BLUE = "4472C4"
NAVY = "102A43"
ORANGE = "ED7D31"
LIGHT_BLUE = "EAF1FB"
LIGHT_ORANGE = "FDF0E7"
LIGHT = "F5F7FA"
MID = "5B6777"
DARK = "263238"
WHITE = "FFFFFF"
S_COLOR = "ED7D31"
T_COLOR = "4472C4"
FONT = "Microsoft YaHei"

# page, title, body items, source, assets, layout, note
PAGES = [
(1,"Trench 与 SGT VDMOS 总剂量辐照效应对比及 TCAD 建模验证",["⁶⁰Co γ 总剂量辐照","S1–S3 / T1–T3","0–60 krad(Si) 七剂量纵向跟踪","统一提参 · 双器件 TCAD 验证"],"原讲稿版模板；项目实验与仿真证据链",[],"cover","本研究对六只器件开展七剂量纵向实验，比较转移与关断输出，并用统一算法提取参数。进一步以两类器件的 TCAD 结果检验可解释趋势、残差和适用边界。"),
(2,"高辐射环境会持续改变功率 MOS 器件的关键电学参数",["辐射激发电子—空穴对","氧化层与界面缺陷累积","Vth 负移、gm 降低、SS 增大","关断输出与高场响应改变"],"Liu et al., 2008；Wang et al., 2023/2024",["02_背景与结构/ChatGPT Image 2026年7月23日 15_43_15.png"],"image_right","电子—空穴对、有效正电荷和界面态构成通用机制框架，但失效程度取决于器件结构、偏置、剂量和工艺。本研究未测陷阱谱，因此不分离具体缺陷贡献。"),
(3,"TCAD 将长周期辐照实验转化为可迭代、可解释的模型研究",["实验：提供真实器件响应","TCAD：控制参数并观察内部状态","实验回答“发生了什么”","仿真回答“当前模型能解释到什么程度”"],"六器件七剂量实验；T/S 七剂量仿真",["02_背景与结构/ChatGPT Image 2026年7月23日 15_45_34.png"],"image_left","实验是主证据，TCAD 用统一参数同时检验 Vth、gm 和 SS。仿真依赖假设，不能替代标定实验，也不用于宣称自主发现物理规律。"),
(4,"已有研究建立 TID 机制框架，但结构对比仍需具体样品证据",["通用机制：氧化层俘获与界面态","结构比较：Trench 与 SGT 的响应可能不同","本项目：同器件、同算法、同证据等级","文献幅度不直接移植到本批器件"],"Liu et al., 2008；Wang et al., 2023/2024",[],"cards","文献覆盖机制、结构比较与 TCAD 分析，但器件、偏置、剂量和提参方法不同。本项目以本批样品的纵向实验和标定结果为准。"),
(5,"三篇 TID 文献分别支撑机制、结构比较与 TCAD 分析",["Liu 2008｜商业 Trench 的 Vth 与直流退化","Wang 2023｜100 V SGT/Trench 对比","Wang 2024｜SGT 的剂量、偏置与 TCAD","三层证据：通用机制 → 结构比较 → SGT 深入分析"],"Liu et al., IEEE TNS, 2008；Wang et al., 2023/2024",[],"cards","三篇文献构成从通用机制到结构比较，再到 SGT 深入分析的证据链。文献中的击穿变化不代表本项目已完成 BV 拟合。"),
(6,"研究闭环从真实测量延伸到两类器件的可追溯 TCAD 验证",["实验测量","统一处理","TCAD 建模","实验—仿真比较","证据审计"],"项目研究闭环与正式交付证据链",["02_背景与结构/route.png"],"full_image","闭环不仅比较拟合成功的部分，也保留残差、失败和覆盖不足。正式结论可回溯到 deck、PLT、日志、CSV、运行状态和哈希。"),
(7,"六只器件在七个剂量点进行同器件纵向跟踪",["S1–S3：SGT VDMOS","T1–T3：Trench VDMOS","0 / 10 / 20 / 30 / 40 / 50 / 60 krad(Si)","20 V 辐照偏置；198 s / 10 krad","转移：VDS=0.1 V；输出：VGS=−20 V"],"ST_TID_综合成果报告；实验记录",["02_背景与结构/site_irradiator.jpg","02_背景与结构/site_room.jpg","02_背景与结构/site_source.jpg"],"timeline","每组只有三只器件，七个剂量点是同一器件的纵向跟踪，不是独立样本扩增。高压端约 10 mA 为仪器限流，不作为本征剂量指标。"),
(8,"Trench 与 SGT 通过不同沟槽电极和氧化层布局控制内部电场",["Trench：沟槽栅与侧壁沟道","SGT：控制栅、屏蔽栅与厚氧化层","多氧化层使 TID 耦合更复杂","不预设某一结构在全部指标上更优"],"项目结构图；Wang et al., 2023/2024",["02_背景与结构/T.png","02_背景与结构/S.png"],"two_images","SGT 的场板与屏蔽设计可改善常规性能，但不预设其全部指标更耐辐照。文献结构尺寸和掺杂不套用于本项目器件。"),
(9,"统一算法从原始曲线重新提取 Vth、gm、SS 与输出指标",["28 个工作簿 / 97 张表","42 条转移记录 / 55 条输出记录","Vth：11 点二阶 Savitzky–Golay 最大 gm 切线","SS：21 点滑窗，R²≥0.98","异常、扫描范围和质量限制全部保留"],"analysis_config.json；数据质量报告",[],"process","原表中的 GM/VT 不能直接复用，因此使用统一算法重新计算。这样确保两组器件使用同一把尺子，并保留异常和窗口限制。"),
(10,"六只器件共同显示阈值负移与开启能力下降",["60 krad：S 组 ΔVth=−6.120±0.045 V","60 krad：T 组 ΔVth=−9.969±0.337 V","两组 gm 均下降约 60%","T 组 SS 恶化更显著"],"summary_facts.json；transfer_summary.csv",["03_实验结果/可编辑图表/transfer_six_device_gallery.png"],"full_image","共同趋势不是由单只器件离群造成。组均值用于总览，逐器件页面继续保留离散与异常，不能把曲线简化为刚性平移。"),
(11,"S1：七剂量转移曲线显示稳定负移与开启区斜率下降",["ΔVth=−6.088 V","gm 变化 −59.1%","SS=328.8 mV/dec","阈值移动与曲线形态变化同时存在"],"summary_facts.json；S1 主扫描",["03_实验结果/S_T_VDMOS_TID_完整成果展示版_讲稿版_10.png"],"legacy","S1 同时出现阈值负移、开启区斜率下降和亚阈形态变化，因此不是简单的刚性平移。"),
(12,"S2：阈值负移与 S1 接近，组内一致性较好",["ΔVth=−6.172 V","gm 变化 −59.2%","SS=288.3 mV/dec","约 6 V 负移不是单器件现象"],"summary_facts.json；S2 主扫描",["03_实验结果/S_T_VDMOS_TID_完整成果展示版_讲稿版_11.png"],"legacy","S2 的 Vth 和 gm 与 S1 高度接近，SS 略低但仍明显劣于初始状态。"),
(13,"S3：跨导下降略大，但阈值负移仍与 S 组高度一致",["ΔVth=−6.099 V","gm 变化 −62.5%","SS=348.7 mV/dec","Vth 集中，gm 与 SS 存在器件离散"],"summary_facts.json；S3 主扫描",["03_实验结果/S_T_VDMOS_TID_完整成果展示版_讲稿版_12.png"],"legacy","S3 的阈值负移与 S1、S2 接近，但开启能力和亚阈区退化略强，说明不同指标的组内离散程度并不相同。"),
(14,"T1：阈值负移超过 10 V，亚阈区明显展开",["ΔVth=−10.339 V","gm 变化 −56.6%","SS=1698.3 mV/dec","Vth 与 SS 退化显著强于 S 组"],"summary_facts.json；T1 主扫描",["03_实验结果/S_T_VDMOS_TID_完整成果展示版_讲稿版_13.png"],"legacy","T1 的 Vth 和 SS 退化明显，但 gm 相对降幅并未超过全部 S 器件，因此三个指标必须分别解读。"),
(15,"T2：跨导下降最大，阈值与 SS 保持 T 组共同趋势",["ΔVth=−9.885 V","gm 变化 −64.4%","SS=1690.1 mV/dec","六只器件中 gm 相对降幅最大"],"summary_facts.json；T2 主扫描",["03_实验结果/S_T_VDMOS_TID_完整成果展示版_讲稿版_14.png"],"legacy","T2 的 gm 降幅最大，但 Vth 和高剂量 SS 与 T1、T3 仍保持共同趋势。"),
(16,"T3：三只 T 器件均进入约 1.6–1.7 V/dec 的高 SS 区间",["ΔVth=−9.682 V","gm 变化 −62.7%","SS=1633.4 mV/dec","30–60 krad 的高 SS 不是偶发窗口"],"summary_facts.json；T3 主扫描",["03_实验结果/S_T_VDMOS_TID_完整成果展示版_讲稿版_15.png"],"legacy","T3 再次表明 T 组在中高剂量进入共同的高 SS 区间，这一趋势不是单器件或单窗口偶发。"),
(17,"两组 gm 均下降约 60%，但 Vth 与 SS 的退化幅度不同",["S/T ΔVth：−6.120±0.045 / −9.969±0.337 V","SS 增量：133.6 / 1448.6 mV/dec","T/S 的 SS 增量比约 10.8×","gm 是共同退化，SS 是最大结构差异"],"transfer_summary.csv；transfer_parameters.csv",["03_实验结果/S_T_VDMOS_TID_完整成果展示版_讲稿版_16.png"],"legacy","T 组比 S 组多负移约 3.85 V。gm 是两组共同退化维度，而 SS 构成最显著的组间差异。"),
(18,"S 组转移曲线随剂量稳定负移，三只器件保持高度一致",["Vth：2.073±0.014 → −4.047±0.035 V","ΔVth=−6.120±0.045 V","gm 下降 60.25±1.96%","SS：188.3 → 321.9 mV/dec"],"实验统计与统一提参结果",["03_实验结果/transfer_semilog_S.png","03_实验结果/transfer_linear_S.png"],"two_images","S 组同时出现横向平移、开启区和亚阈形态变化。相对 T 组，S 组的 SS 恶化较缓，但这些指标并非单一缺陷的直接读数。"),
(19,"S 组七剂量指标给出连续、可复核的退化轨迹",["Vth：2.073 → −4.047 V","gm：0 → −60.25%","SS：188.3 → 321.9 mV/dec","40 krad SS 局部回落","Id@30 V 与 V@1µA 呈非单调细节"],"ST_TID_综合成果报告，表 2-1",[],"trend","Vth 与 gm 总体持续退化，但 40 krad 的 SS 局部回落、输出指标非单调。不能把全部指标压缩成同一个剂量函数。"),
(20,"T 组阈值负移接近 10 V，亚阈区在中高剂量快速恶化",["Vth：3.143±0.032 → −6.826±0.350 V","ΔVth=−9.969±0.337 V","gm 下降 61.21±4.10%","SS：225.4 → 1674.0 mV/dec"],"实验统计与统一提参结果",["03_实验结果/transfer_semilog_T.png","03_实验结果/transfer_linear_T.png"],"two_images","20 krad 时 T 组平均 Vth 已转为负值，30–60 krad 进入共同高 SS 区。SS 恶化提示栅控退化，但不能直接分离缺陷分量。"),
(21,"T 组 Vth 与 SS 的退化快于输出漏电变化",["Vth：3.143 → −6.826 V","SS：225.4 → 1674.0 mV/dec","Id@30 V：1.22×10⁻¹⁰ → 1.95×10⁻¹⁰ A","V@1µA：92.7 → 108.0 V","转移与关断输出并不同步"],"ST_TID_综合成果报告，表 3-1",[],"trend","T 组转移参数强烈退化，但低压关断漏电和高压代理变化较小。不能用一个总体耐辐照结论替代分指标解读。"),
(22,"两组 gm 均下降约 60%，开启能力衰减是共同退化维度",["SGT：−60.25±1.96%","Trench：−61.21±4.10%","两组均值仅相差约 1 个百分点","相对变化接近不代表绝对 gm 相同"],"实验统计",["03_实验结果/gm_change.png"],"image_right","迁移率、界面散射、曲线形态与串联电阻都可能参与 gm 变化。由于没有独立迁移率或陷阱谱测量，不能唯一归因。"),
(23,"T 组 SS 增量约为 S 组的 10.8 倍，亚阈区差异最显著",["S：188.3 → 321.9 mV/dec","T：225.4 → 1674.0 mV/dec","增量比约 10.8×","该比值是增量比，不是绝对值比"],"实验统计；统一 SS 滑窗算法",["03_实验结果/subthreshold_swing.png"],"image_right","T 组低中剂量 SS 离散仍需保留，但高剂量共同趋势明确。SS 不能直接反演界面态密度。"),
(24,"SGT 阈值负移更小，Trench 在 20 krad 后平均 Vth 已转为负值",["S：2.073 → −4.047 V","T：3.143 → −6.826 V","T 组多负移约 3.85 V","T 组 20 krad：−0.484 V"],"实验统计",["03_实验结果/threshold_shift.png","03_实验结果/threshold_voltage.png"],"two_images","本批 SGT 的阈值稳定性更好，但两组不是同工艺仅改变分裂栅的严格对照，因此不能把差异唯一归因于分裂栅。"),
(25,"本批结果提示 SGT 的电极与氧化层布局提高了转移特性稳定性",["实验事实：|ΔVth| 6.120 vs 9.969 V","实验事实：SS 增量 133.6 vs 1448.6 mV/dec","共同事实：gm 均下降约 60%","候选解释：电极、氧化层与内部电场共同影响"],"实验结果；Wang et al., 2023/2024",["02_背景与结构/T.png","02_背景与结构/S.png"],"fact_hypothesis","结构解释是由实验差异和结构知识提出的假设，需要 TCAD 和更严格的同工艺对照继续验证，不能推广为所有 SGT 的普适结论。"),
(26,"输出曲线需区分低压漏电、高压代理指标与仪器限流",["Id@30 V：低压关断漏电","V@1µA：高压首达代理指标","约 10 mA 平台：仪器限流","V@1µA 不直接等同标准 BV"],"实验指标定义",[],"diagram","输出测试固定 VGS=−20 V，只有单一栅压扫描，不是多栅压输出族。三个区段描述不同现象，不能混用。"),
(27,"六只器件的输出响应呈现明显组间差异与非单调细节",["S 组：低压漏电更大，部分首扫出现尖峰","T 组：低压变化较小，V@1µA 约 100 V","中间剂量与异常点全部保留","输出比 Vth 更难用单参数解释"],"六器件七剂量首扫输出",["03_实验结果/可编辑图表/output_six_device_gallery.png"],"full_image","输出受关断裕量、漏电路径、局部电场、陷阱占据和测量历史共同影响，因此必须保留非单调和中间剂量细节。"),
(28,"S1：60 krad 的 30 V 漏电增幅为 1749%，高压代理降至 36 V",["Id@30 V：+1749.3%","V@1µA：36.0 V","中间剂量并非严格单调","10 mA 平台不是本征击穿"],"S1 输出主扫描",["03_实验结果/S_T_VDMOS_TID_完整成果展示版_讲稿版_19.png"],"legacy","S1 的端点变化明显，但中间剂量并不严格单调。高压限流平台只反映仪器合规状态，不代表本征击穿。"),
(29,"S2：三只 S 器件中 30 V 漏电相对增幅最大",["Id@30 V：+2874.4%","V@1µA：35.0 V","极低基线会放大相对百分比","需同时查看绝对电流与完整曲线"],"S2 输出主扫描",["03_实验结果/S_T_VDMOS_TID_完整成果展示版_讲稿版_20.png"],"legacy","S2 的相对增幅最大，但百分比受到极低初始电流放大，不能仅凭百分比判定失效等级。"),
(30,"S3：相对增幅较低，但高压代理同样降至约 35 V",["Id@30 V：+897.1%","V@1µA：35.0 V","低压漏电存在组内离散","高压首达集中在 35–36 V"],"S3 输出主扫描",["03_实验结果/S_T_VDMOS_TID_完整成果展示版_讲稿版_21.png"],"legacy","S3 的低压漏电增幅低于 S1、S2，但高压首达代理仍集中在同一区间。两个指标对应曲线的不同区段。"),
(31,"T1：低压漏电仅小幅增加，高压代理保持在 108 V",["Id@30 V：+72.7%","V@1µA：108.0 V","转移退化强，但关断输出变化弱","漂移区与关断条件需单独分析"],"T1 输出主扫描",["03_实验结果/S_T_VDMOS_TID_完整成果展示版_讲稿版_22.png"],"legacy","T1 的转移参数退化很强，但强制负栅关断下输出变化较弱，说明转移与高压区不能用同一结论替代。"),
(32,"T2：30 V 漏电增幅为 58.9%，高压代理为 109 V",["Id@30 V：+58.9%","V@1µA：109.0 V","端点与 T1 接近","中间剂量仍保留非单调细节"],"T2 输出主扫描",["03_实验结果/S_T_VDMOS_TID_完整成果展示版_讲稿版_23.png"],"legacy","T2 的端点结果与 T1 接近，但中间剂量仍可能非单调，因此正式结论基于完整曲线而非只看端点。"),
(33,"T3：三只 T 器件的 60 krad 输出指标集中在同一区间",["Id@30 V：+48.9%","V@1µA：107.0 V","T1–T3 高压代理集中在 107–109 V","与 Vth/SS 剧烈退化形成反差"],"T3 输出主扫描",["03_实验结果/S_T_VDMOS_TID_完整成果展示版_讲稿版_24.png"],"legacy","T 组输出端点较集中，而 Vth 和 SS 退化显著，进一步说明不同工作区的敏感机制并不相同。"),
(34,"S 组首扫尖峰在复扫中消失，表现出扫描历史相关响应",["首扫约 30 V 出现尖峰或高漏电","复扫电流显著降低","现象覆盖多条配对曲线","机制尚未通过扫速、间隔和温度控制唯一分离"],"13 条 S 组首扫—复扫配对",["03_实验结果/output_semilog_S.png"],"image_right","漏电降低是重复出现的测量事实，但未控制等待时间、扫速和温度。只称扫描历史相关响应或恢复现象，不把陷阱释放写成唯一机制。"),
(35,"S/T 输出指标给出相反的高压代理变化方向",["S：Id@30 V +1840±992%","S：V@1µA 106.3±6.4 → 35.3±0.6 V","T：Id@30 V +60.2±11.9%","T：V@1µA 92.7±1.2 → 108.0±1.0 V"],"output_summary.csv；output_parameters.csv",["03_实验结果/S_T_VDMOS_TID_完整成果展示版_讲稿版_25.png"],"legacy","两个指标观察曲线的不同区段，方向不同不构成数据矛盾。S 组百分比还受到极低初始漏电的放大。"),
(36,"13 条 S 组复扫的 30 V 漏电全部降低 70.65%–96.19%",["13 / 13 条配对曲线均降低","降低幅度：70.65%–96.19%","13 条是 paired traces，不是 n=13","替代解释包括自热、接触和仪器状态"],"recovery_parameters.csv；相关恢复文献",["03_实验结果/S_T_VDMOS_TID_完整成果展示版_讲稿版_26.png"],"legacy","方向重复性很强，且与陷阱释放、复合或重分布相容。同时必须保留自热、接触和仪器状态等替代解释。"),
(37,"S3 60 krad 的恢复覆盖较宽电压范围，而非单点变化",["复扫曲线整体低于首扫","40–70 V 局部拐点和电流同时改变","30 V 只是代表采样点","个案用于补强形态证据"],"S3 60 krad 首扫—复扫配对",["03_实验结果/S_T_VDMOS_TID_完整成果展示版_讲稿版_27.png"],"legacy","该案例说明变化覆盖较宽电压范围，但不能替代 13 条配对统计，也不能独自识别陷阱类型。"),
(38,"高场扫描可能通过陷阱电荷释放与重分布形成表观恢复",["TID 形成俘获电荷","首扫高场改变占据状态","释放、复合或迁移改变有效电荷","复扫漏电降低","验证缺口：扫速、间隔、温度与方向矩阵"],"Oldham & McLean；Schwank；Lelis；Gao",[],"process","文献支持可逆退火与去俘获背景，但本实验没有完成机制分离。这里展示的是与现象相容的候选链路，而不是已唯一证明的微观机制。"),
(39,"SGT 与 Trench 的输出响应方向相反，结构敏感区并不相同",["SGT：漏电大增、V@1µA 降至约 35 V、历史相关","Trench：漏电小增、V@1µA 约 100 V、端点集中","SGT 转移更稳，但输出更敏感","不能形成单一的总体耐辐照排名"],"实验输出统计；固定 VGS=−20 V",["03_实验结果/output_semilog_S.png","03_实验结果/output_semilog_T.png"],"two_images","SGT 和 Trench 在不同工作区表现出相反排序，说明敏感区并不相同。结论限于本批样品和当前负栅关断条件。"),
(40,"SGT 的厚氧化层与屏蔽电极可能增强辐照电荷对漂移区电场的影响",["厚氧化层可能积累更多有效电荷","屏蔽场板可能改变空间电势耦合","局部电场扰动可能激活漏电路径","结构解释仍需 TCAD 与严格对照验证"],"实验结果；Wang et al., 2023/2024",["02_背景与结构/S.png","02_背景与结构/T.png"],"fact_hypothesis","这一结构解释与实验差异和公开文献相容，但没有局部场测量或完全同工艺对照，因此只作为待检验假设。"),
(41,"T 器件结构、网格和基础参数固定，剂量只改变陷阱参数",["二维 Trench，300 K","P-body：1.33×10¹⁷ cm⁻³","AreaFactor：30335","转移 VDS=0.1 V","几何、掺杂、面积和网格保持不变"],"tid_campaign.json；T 仿真技术报告",["04_T器件仿真/simulation/structure_mesh.png"],"image_left","固定结构、掺杂、面积和陷阱形状，避免用结构漂移吸收剂量效应。当前结论受二维、300 K 和偏置范围限制。"),
(42,"低参数 Not/Nit 剂量函数将七个剂量点映射到同一模型",["Not=3.0×10¹⁰+4.96×10¹²(D/60)^0.87 cm⁻²","Nit=1.4×10¹¹+4.0×10¹²(D/60)^0.70 cm⁻²","D 的单位为 krad(Si)","七点共同标定，不逐剂量独立调参"],"tid_campaign.json；T 仿真报告",["04_T器件仿真/dose_trap_mapping.svg"],"image_right","低参数单调函数避免逐剂量独立调参。Not 和 Nit 是当前范围的等效参数化模型，不是直接测得的陷阱密度或普适定律。"),
(43,"Not 主导阈值静电位移，Nit 用于约束亚阈区与栅控退化",["Not → 表面势与 Vth","Nit → SS、gm 与界面散射","两类参数构成可辨识压缩模型","改善 SS 却过度压低 gm 时，应承认模型表达不足"],"MOS TID 通用机制；T 参数扫描",[],"diagram","参数职责来自模型构造，不代表缺陷分量已被直接测得或唯一分离。继续无约束增加参数会造成复杂度扩散。"),
(44,"T 组七剂量转移仿真复现阈值负移主趋势",["实验与仿真按同剂量同色展示","半对数观察 Vth 与 SS","线性坐标观察开启区与 gm","趋势吻合不等于整曲线完全重合"],"真实 VM SVisual；transfer_d00–d60.csv",["04_T器件仿真/comparison/comparison_transfer_semilog_T.png"],"full_image","七剂量共同标定复现主趋势，但绝对电流和局部形态仍有残差。没有通过删点、逐剂量缩放或改变 AreaFactor 隐藏残差。"),
(45,"T 模型对 ΔVth 与 gm 达到稳定的范围内定量一致性",["ΔVth RMSE=0.0771 V","60 krad ΔVth 误差=0.00177 V","gm 变化 RMSE=5.97 pct-pt","60 krad gm：仿真 −63.22%，实验 −61.21%"],"simulation_metrics.csv；fit_error_table.csv",["04_T器件仿真/t_simulation_metric_compare.svg"],"image_right","这些误差描述共同标定范围内的一致性，不能称为盲测预测精度。"),
(46,"SS 残差揭示均匀陷阱模型无法同时解释全部退化",["60 krad SS：仿真 892 mV/dec","60 krad SS：实验 1674 mV/dec","增加均匀 Acceptor 可改善 SS","但会过度压低 gm，形成明确权衡"],"T 仿真报告；A–I 参数扫描",[],"tradeoff","该负结果必须保留。可能还需要非均匀陷阱谱、迁移率、寄生漏电或三维终端效应，当前只能标记为限定支持。"),
(47,"无雪崩预击穿输出只能用于 30 V 内的限定趋势比较",["VGS=−20 V，无雪崩反馈，VDS≤30 V","D30：15.6–30 V / 17 点 / partial","D40：0.003–30 V / 163 点 / 6 个低压负点","仿真 V@1µA 未获得"],"current_output_provenance.json",["04_T器件仿真/comparison/comparison_output_semilog_T.png"],"image_right","缺段不补点，低压负点不删除。Id@30 V 不能替代 V@1µA，这一页不是 BV 拟合。"),
(48,"实验与 SDevice 曲线叠加保留离散、缺段和异常",["21 条 T 实测曲线 + 7 条 SDevice 曲线","颜色表示剂量，线型表示器件或仿真","不均值化、不重采样、不删点","D30 partial 状态明确标注"],"origin_overlay_manifest.json；origin_qa.json",["04_T器件仿真/comparison/comparison_id30_T.png"],"image_right","半对数展示使用可追溯的 abs(Id) 派生列。Origin 只负责同图展示，不替代 PLT、日志和 SVisual 证据。"),
(49,"高压求解在约 30 V 后显著变慢，失败尝试构成模型边界",["Poisson 与耦合稳态初始化","分段偏压 ramp 与保守步长","保存失败 deck、日志、终止电压和租约","数值失败不等于物理击穿"],"高场失败证据与运行记录",["04_T器件仿真/comparison/comparison_high_voltage_context_T.png"],"image_left","最后收敛点不能称为 BV。失败记录用于界定模型的可计算范围，而不是被包装成物理结果。"),
(50,"IIC 是条件化高场判据，不能替代 V@1µA 或标准 BV",["IIC=1.1 作为触发阈值","到上限未触发：仅形成下界","触发前失败：无 IIC/BV 结果","IIC、V@1µA、Id@30 V 与预击穿输出不可互换"],"highfield_closure.json；T 仿真报告",[],"status","IIC 只对当前模型、偏压路径和合格收敛状态有效。仿真 V@1µA 仍是阻塞项。"),
(51,"终态电场图支持空间形态比较，但不能单独证明 BV 不变",["E=−∇φ","比较高场区位置、形态与剖面","0/60 krad 高场位置接近","终止电压与 IIC 状态不同，峰值不可无条件比较"],"真实 VM SVisual；release_manifest.json",["04_T器件仿真/simulation/electric_field_0krad.png","04_T器件仿真/simulation/electric_field_60krad.png"],"two_images","电场图只支持未观察到明显空间迁移这一限定结论。高场位置接近不等于 BV 不变，也不能代替标准击穿判据。"),
(52,"T 输出实验与仿真高场形态只形成定性对应",["实验 V@1µA：92.7±1.2 → 108.0±1.0 V","仿真高场位置总体接近","实验输出端相对稳定","不同口径只允许定性趋势对应"],"实验统计；IIC 与电场证据",["03_实验结果/output_semilog_T.png","04_T器件仿真/simulation/electric_field_60krad.png"],"two_images","V@1µA、IIC 和终态电场是不同口径。本页不宣称 BV 已完成定量拟合。"),
(53,"T 仿真结论同时包含拟合成果、模型残差和数值边界",["已支持：ΔVth RMSE 0.0771 V","已支持：gm RMSE 5.97 pct-pt","受限：SS 残差、D30 缺段、D40 负点","阻塞：仿真 V@1µA 与无条件 BV 结论"],"T_device_TID_simulation_report.md；交付 manifest",[],"supported_limited","模型的可信范围由成功结果、残差、失败和证据链共同定义，而不是只展示拟合较好的部分。"),
(54,"S 型七剂量转移曲线已完成实验—仿真同图比较",["0–60 krad 七个剂量点","实验虚线、仿真实线，同剂量同色","中低剂量部分趋势接近","高剂量 Vth 负移仍未充分复现"],"S 仿真七剂量比较图与映射文件",["05_S器件仿真/七剂量转移/figures/transfer_semilog.svg","05_S器件仿真/七剂量转移/figures/transfer_linear.svg"],"two_images","S 仿真与 T 仿真采用相同展示口径：曲线、指标、误差和边界。当前结果用于同时呈现接近部分和明显偏差，不称拟合完成。"),
(55,"S 模型能描述 gm 下降，但高剂量 Vth 负移仍明显不足",["Vth 0 krad：仿真/实验 1.090 / 2.029 V","Vth 60 krad：−0.051 / −4.066 V","gm 0 krad：0.07091 / 0.12512 S","gm 60 krad：0.05274 / 0.05682 S"],"metrics_comparison.csv",["05_S器件仿真/七剂量转移/figures/vth_vs_dose.svg","05_S器件仿真/七剂量转移/figures/gm_vs_dose.svg"],"two_images","gm 高剂量趋势较接近，但基线偏低；Vth 在 20 krad 偶然接近，高剂量偏差约 4 V。七点 Vth RMSE 约 2.42 V 仅描述当前差异。"),
(56,"严格 SS 算法下仅 20 krad 仿真曲线具备有效比较窗口",["21 点滑窗，R²≥0.98","仅 20 krad 通过质量判据","仿真/实验：285.34 / 207.63 mV/dec","其余 6 个剂量不填零、不插值、不连线"],"S仿真 qa.json；metrics_comparison.csv",["05_S器件仿真/七剂量转移/figures/ss_vs_dose.svg"],"image_right","当前不支持 S 七剂量 SS 拟合。下一步应补充低电流区扫描与点密度，而不是放宽算法标准。"),
(57,"S 型低漏压终点的局部电场峰值仅变化约 0.18%",["固定 Y=−6.0 µm","0 krad：9.654×10⁵ V/cm","60 krad：9.671×10⁵ V/cm","变化 +0.1765%，峰位约 X=0.803 µm"],"peak_metrics.csv；svisual_manifest.json",["05_S器件仿真/低压电场/figures/field_x_cutline_overlay.svg","05_S器件仿真/低压电场/figures/field_peak_vs_dose.svg"],"two_images","这是 VDS≈0.1 V、VGS≈5 V 的转移终点，只说明局部场小幅重分布。它不是关断高压或击穿条件，不能用于推断 BV。"),
(58,"S 模型应按基线、剂量映射和亚阈扫描的顺序迭代",["第一层：修正 0 krad 的 Vth 与 gm 基线","第二层：重建中高剂量 Not/Nit 映射","第三层：补充低电流区与 SS 扫描密度","分层迭代，避免参数复杂度扩散"],"S Vth/gm/SS 比较与 QA",[],"process","当前问题不是再调一个参数即可解决。应先校准基线，再调整剂量映射，最后补齐亚阈区扫描，否则不同误差会互相掩盖。"),
(59,"S 型仿真已形成真实证据链，但模型验证仍处于迭代阶段",["已形成：七剂量叠图与 Vth/gm 比较","已形成：20 krad SS 与低压电场","待迭代：高剂量 Vth、0 krad gm、六剂量 SS","未覆盖：S 输出与高压验证"],"S仿真 manifest、mapping 与 QA",["05_S器件仿真/七剂量转移/figures/svisual_transfer_semilog_0to60K_clean.png"],"supported_limited","S 仿真的价值在于定位哪些剂量和指标接近、哪些尚未解决。GUI 导出用于证明图像链，不替代定量数据和 QA。"),
(60,"本研究形成器件对比、恢复现象、双模型与 AI 流程四类创新",["结构对比｜同器件七剂量实验","恢复现象｜13/13 配对曲线降低","双模型｜T 定量标定 + S 阶段诊断","工程流程｜受监督、可追溯的 AI 协作"],"本项目实验、仿真与监督工程材料",[],"cards","创新不是单条曲线拟合，而是结构差异实验、恢复现象、双模型验证与工程证据流程组成的体系。"),
(61,"TCAD 的主要瓶颈是反复收敛、比较与证据筛选",["参数之间存在耦合","数值失败与高场耗时难以避免","deck、mesh、log、PLT、偏压和哈希需逐项核验","AI 用于比较、监控和证据整理"],"T/S 参数扫描、失败证据与监督规范",[],"diagram","AI 降低跨文件搜索、监控和核验成本；研究者仍负责物理边界、参数批准、停止条件和结果解释。"),
(62,"AI 协作采用“探索—审批—执行—监控—复核”的受监督状态机",["探索：多角色只读分析","审批：研究者确认候选与边界","执行：白名单单候选运行","监控：只读报告，不自动改参","复核：验证、SVisual 与人工晋级"],"监督状态机；人工审批与派发规则",[],"process","AI 不自主调参。新增参数、重试、扩大范围或晋级正式结果都要经过研究者确认。"),
(63,"单核原子租约让多任务并发保持可控且失败关闭",["所有 SDevice 使用 --threads 1","CPU0 保留，最多 4 个租约槽位","原子 claim → 亲和性验证 → 释放并留证","依赖链串行，仅独立剂量点允许并发"],"sentaurus-core-policy.json；sdevice_core_lease.sh",[],"process","租约、亲和性或哈希异常时立即拒绝执行。只有共同网格已完成且没有依赖的剂量点才能并发。"),
(64,"异步监控只报告状态，不自动重试或修改参数",["30 s → 1 min → 2 min 的渐进观察","默认 45/60 s 轮询","检查 manifest、远端进程与 atomic lease","异常只上报，不改参、不重试、不终止"],"监控脚本；监督状态机",[],"timeline","监控角色保持只读，完成或异常均交回主代理和研究者决定下一步，避免自动化越权。"),
(65,"AI 与研究者职责分离，避免自动化越过物理与证据边界",["AI：搜索、汇总、候选、监控、CSV/哈希核验","研究者：器件事实、物理合理性、参数与预算审批","研究者：停止条件、结果晋级与答辩表述","AI 辅助，不替代责任主体"],"人工审批与证据复核规则",[],"two_columns","AI 提升工程协作效率，但不拥有物理结论和发布结果的权力。"),
(66,"失败、部分结果和异常点均保留为可追溯证据",["D30 partial 不补点","D40 保留 6 个低压负点","高场失败不包装成 BV","S 高剂量 Vth 偏差与 6/7 SS 无窗如实保留"],"监督证据规则；T/S QA",[],"cards","这一流程防止缺段拼接、异常静默删除和数值收敛被误标为物理结果。"),
(67,"自动化流程已管理 22,627 个原始点和多轮仿真证据",["14,322 个转移原始点","8,305 个输出原始点","42 条转移记录","55 条输出记录","13 条复扫配对","T/S 各 7 个剂量仿真点"],"数据质量报告；versioned campaigns",[],"kpi","这里用真实工作规模说明 AI 的作用，不虚构效率提升倍数，也不把原始点数量当作统计样本量。"),
(68,"真实虚拟机工程与 AI 操作记录展示了持续迭代过程",["分析与任务编排","运行监控与错误识别","Workbench / SVisual 工程","分析 → 下发 → 监控 → 复核 → 调整"],"用户提供的真实操作截图",["06_AI辅助流程/AI操作中图2.png","06_AI辅助流程/AI操作中图像1.png","06_AI辅助流程/虚拟机工程.png"],"three_images","这些截图证明 AI 参与了分析、下发、监控和复核链路。每轮参数调整和结果有效性仍由研究者确认。"),
(69,"四项创新将实验现象推进为可解释、可追溯的模型成果",["结构比较：转移与输出排序不同","恢复现象：13/13 配对降低","双模型：T 定量标定，S 阶段诊断","AI 流程：审批、租约、监控与证据归档"],"前述实验、仿真与流程证据",[],"process","四项创新形成递进关系：从提出结构差异问题，到发现历史相关响应，再到模型检验和工程支撑。"),
(70,"两类器件均发生强烈转移退化，但输出响应方向不同",["SGT｜ΔVth −6.120±0.045 V｜gm −60.25±1.96%","SGT｜SS 321.9±30.7 mV/dec｜V@1µA 35.3 V","Trench｜ΔVth −9.969±0.337 V｜gm −61.21±4.10%","Trench｜SS 1674.0±35.4 mV/dec｜V@1µA 108.0 V"],"实验统计",[],"matrix","两类器件均有强烈转移退化，但 SGT 的转移更稳定、Trench 的关断输出更稳定。选型必须按工作区和失效模式，而不是总体排名。"),
(71,"T 模型实现定量标定，S 模型完成阶段性验证与问题定位",["T：ΔVth RMSE 0.0771 V","T：gm RMSE 5.97 pct-pt，保留 SS 与高场边界","S：完成七剂量转移与低压电场比较","S：高剂量 Vth、SS、输出与高压仍需迭代"],"T/S 仿真证据与 QA",[],"two_columns","两类模型使用相同展示方式，但成熟度不同。T 是共同标定范围内的定量模型，S 是阶段性验证与诊断结果。"),
(72,"关键证据边界集中保留，避免成果展示产生误解",["每组 n=3，七剂量是纵向跟踪","V@1µA 是高压代理，不是标准 BV","S 恢复机制尚未唯一分离","T 七点是共同标定，不是盲测预测","T SS/D30/D40/V@1µA 与 S 高剂量结果均有边界","AI 不代替研究责任主体"],"全套实验、仿真与监督证据",[],"cards","这些边界不会削弱结论，反而明确了结论能回答什么、不能回答什么，避免被误读为普适规律或全部完成。"),
(73,"下一步工作将围绕 S 模型、高压验证与恢复机制展开",["S 模型：校准基线与高剂量剂量映射","高压闭环：统一 V@1µA、IIC 与 BV 口径","恢复机制：等待时间、扫速、温度与方向矩阵","继续采用受监督、可追溯的 AI 工程流程"],"基于当前 T/S 模型边界",[],"process","下一阶段分为三条主线：完善 S 模型、建立高压闭环、开展恢复机制控制实验，使解释由相容逐步走向可辨识。"),
(74,"谢谢各位老师，请批评指正",["实验数据完整","模型边界可见","仿真证据可追溯","问题与下一步明确"],"本项目答辩材料",["02_背景与结构/route.png"],"closing","汇报结束，感谢各位老师的指导，请批评指正。"),
]


def rgb(hexstr: str) -> RGBColor:
    return RGBColor.from_string(hexstr)


def add_text(slide, text, x, y, w, h, size=16, color=DARK, bold=False,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, margin=0.06):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(margin)
    tf.margin_top = tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = text
    for run in p.runs:
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
    return shape


def add_rect(slide, x, y, w, h, fill, line=None, radius=True):
    typ = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    sh = slide.shapes.add_shape(typ, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = rgb(fill)
    sh.line.color.rgb = rgb(line or fill)
    return sh


def add_header(slide, page, title, section):
    add_rect(slide, 0, 0, 10, 0.13, BLUE, radius=False)
    add_text(slide, section, 0.48, 0.22, 2.0, 0.28, 10.5, BLUE, True)
    add_text(slide, title, 0.48, 0.53, 8.75, 0.64, 25 if len(title) < 34 else 22, NAVY, True)
    add_text(slide, f"{page:02d} / 74", 8.92, 0.22, 0.62, 0.26, 10.5, MID, True, PP_ALIGN.RIGHT)
    add_rect(slide, 0.48, 1.21, 0.78, 0.055, ORANGE, radius=False)


def add_footer(slide, source):
    add_text(slide, f"来源：{source}", 0.48, 5.28, 8.65, 0.22, 9.5, MID)
    add_text(slide, "T = Trench VDMOS　|　S = SGT VDMOS", 7.15, 5.28, 2.37, 0.22, 8.5, MID, False, PP_ALIGN.RIGHT)


def render_vector(path: Path) -> Path:
    CACHE.mkdir(parents=True, exist_ok=True)
    key = hashlib.sha256(str(path).encode("utf-8")).hexdigest()[:16]
    out = CACHE / f"{key}.png"
    if out.exists():
        return out
    ext = path.suffix.lower()
    if ext == ".svg":
        data = path.read_bytes()
        doc = fitz.open("svg", data)
    else:
        doc = fitz.open(path)
    page = doc[0]
    pix = page.get_pixmap(matrix=fitz.Matrix(2.2, 2.2), alpha=False)
    pix.save(out)
    doc.close()
    return out


def resolved_asset(rel: str) -> Path | None:
    path = ROOT / rel
    if not path.exists():
        # known fallbacks
        swaps = {
            "transfer_six_device_gallery.png": "transfer_metrics_overview.png",
            "output_six_device_gallery.png": "output_metrics_overview.png",
            "comparison_transfer_semilog_T.png": "comparison_transfer_metrics_T.png",
        }
        alt = path.with_name(swaps.get(path.name, path.name))
        if alt.exists(): path = alt
        else: return None
    if path.suffix.lower() in {".svg", ".pdf"}:
        return render_vector(path)
    return path


def add_image_fit(slide, path: Path, x, y, w, h, border=True):
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(w / iw, h / ih)
    nw, nh = iw * scale, ih * scale
    lx, ty = x + (w - nw) / 2, y + (h - nh) / 2
    if border:
        add_rect(slide, x, y, w, h, WHITE, "D9E2EC")
    slide.shapes.add_picture(str(path), Inches(lx), Inches(ty), Inches(nw), Inches(nh))


def add_bullets(slide, items, x, y, w, h, color=BLUE, size=16):
    n = max(1, len(items)); gap = 0.1
    each = (h - gap * (n - 1)) / n
    for i, item in enumerate(items):
        yy = y + i * (each + gap)
        add_rect(slide, x, yy, w, each, LIGHT, "D9E2EC")
        add_rect(slide, x, yy, 0.07, each, color, radius=False)
        add_text(slide, item, x + 0.18, yy + 0.05, w - 0.27, each - 0.08, size, DARK, i == 0,
                 valign=MSO_ANCHOR.MIDDLE)


def add_cards(slide, items, x=0.48, y=1.55, w=9.04, h=3.5):
    cols = 2 if len(items) <= 4 else 3
    rows = (len(items) + cols - 1) // cols
    cw = (w - 0.22 * (cols - 1)) / cols
    ch = (h - 0.22 * (rows - 1)) / rows
    for i, item in enumerate(items):
        col, row = i % cols, i // cols
        xx, yy = x + col * (cw + 0.22), y + row * (ch + 0.22)
        fill = LIGHT_BLUE if i % 2 == 0 else LIGHT_ORANGE
        accent = BLUE if i % 2 == 0 else ORANGE
        add_rect(slide, xx, yy, cw, ch, fill, "D9E2EC")
        add_rect(slide, xx, yy, cw, 0.07, accent, radius=False)
        lead, sep, rest = item.partition("｜")
        if sep:
            add_text(slide, lead, xx + 0.18, yy + 0.25, cw - 0.36, 0.42, 17, accent, True)
            add_text(slide, rest, xx + 0.18, yy + 0.72, cw - 0.36, ch - 0.9, 14.5, DARK)
        else:
            add_text(slide, item, xx + 0.18, yy + 0.2, cw - 0.36, ch - 0.35, 15.5, DARK, True,
                     valign=MSO_ANCHOR.MIDDLE)


def add_process(slide, items):
    n = len(items); x0, y, total = 0.48, 2.05, 9.04
    gap = 0.17; w = (total - gap * (n - 1)) / n
    for i, item in enumerate(items):
        x = x0 + i * (w + gap)
        color = BLUE if i % 2 == 0 else ORANGE
        add_rect(slide, x, y, w, 1.55, color)
        add_text(slide, str(i + 1), x + 0.12, y + 0.12, 0.35, 0.35, 18, WHITE, True, PP_ALIGN.CENTER)
        add_text(slide, item, x + 0.13, y + 0.57, w - 0.26, 0.78, 14.5 if n <= 5 else 12.5, WHITE, True,
                 PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        if i < n - 1:
            add_text(slide, "→", x + w, y + 0.47, gap, 0.5, 18, ORANGE, True, PP_ALIGN.CENTER)
    add_text(slide, "全过程受人工审批、质量判据和证据链约束", 1.6, 4.05, 6.8, 0.45, 18, NAVY, True, PP_ALIGN.CENTER)


def section_for(page: int) -> str:
    if page <= 10: return "研究背景与实验设计"
    if page <= 26: return "实验转移特性"
    if page <= 40: return "实验输出与扫描历史"
    if page <= 53: return "T 器件 TCAD 验证"
    if page <= 59: return "S 器件 TCAD 验证"
    if page <= 68: return "受监督 AI 仿真工程流程"
    return "创新、结论与边界"


def add_notes(slide, note: str):
    tf = slide.notes_slide.notes_text_frame
    tf.text = note.strip()


def build():
    OUT.mkdir(parents=True, exist_ok=True)
    prs = Presentation(str(WORK / "base_template.pptx"))
    prs.slide_width = Inches(10); prs.slide_height = Inches(5.625)
    blank = prs.slide_layouts[0]
    assets_log = []
    for page, title, items, source, rel_assets, layout, note in PAGES:
        slide = prs.slides.add_slide(blank)
        # master style background
        bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = rgb(WHITE)
        if layout == "cover":
            add_rect(slide, 0, 0, 10, 5.625, NAVY, radius=False)
            add_rect(slide, 0, 0, 0.18, 5.625, ORANGE, radius=False)
            add_text(slide, title, 0.72, 0.84, 8.65, 1.35, 28, WHITE, True)
            add_rect(slide, 0.72, 2.38, 1.2, 0.08, ORANGE, radius=False)
            add_text(slide, "\n".join(items), 0.72, 2.72, 5.45, 1.55, 17, "D9E2EC")
            add_text(slide, "答辩完整版｜2026.07", 0.72, 4.8, 3.2, 0.35, 12, WHITE, True)
            add_text(slide, "01 / 74", 8.65, 4.8, 0.7, 0.35, 11, "D9E2EC", True, PP_ALIGN.RIGHT)
        elif layout == "closing":
            add_rect(slide, 0, 0, 10, 5.625, NAVY, radius=False)
            p = resolved_asset(rel_assets[0]) if rel_assets else None
            if p:
                add_image_fit(slide, p, 5.5, 0.5, 4.0, 4.2, False)
                pic = slide.shapes[-1]; pic.fill if hasattr(pic, 'fill') else None
            add_text(slide, title, 0.7, 1.55, 5.6, 0.9, 31, WHITE, True)
            add_text(slide, " · ".join(items), 0.72, 2.75, 5.8, 1.0, 17, "D9E2EC")
            add_text(slide, "74 / 74", 8.6, 4.85, 0.7, 0.3, 11, WHITE, True, PP_ALIGN.RIGHT)
        else:
            add_header(slide, page, title, section_for(page))
            paths = []
            for rel in rel_assets:
                p = resolved_asset(rel)
                exists = p is not None
                sha = hashlib.sha256((ROOT / rel).read_bytes()).hexdigest() if (ROOT / rel).exists() else ""
                assets_log.append([page, title, rel, layout, "主选", exists, sha])
                if p: paths.append(p)
            if layout == "legacy" and paths:
                add_image_fit(slide, paths[0], 0.48, 1.43, 7.25, 3.63)
                add_bullets(slide, items[:4], 7.92, 1.43, 1.6, 3.63, ORANGE, 11.5)
            elif layout == "full_image" and paths:
                add_image_fit(slide, paths[0], 0.48, 1.43, 9.04, 3.63)
                add_rect(slide, 0.68, 4.32, 8.64, 0.55, NAVY)
                add_text(slide, items[0] if items else title, 0.88, 4.38, 8.24, 0.39, 14.5, WHITE, True, PP_ALIGN.CENTER)
            elif layout == "two_images" and paths:
                boxes = [(0.48, 1.47, 4.4, 2.9), (5.12, 1.47, 4.4, 2.9)]
                for p, box in zip(paths[:2], boxes): add_image_fit(slide, p, *box)
                add_rect(slide, 0.48, 4.56, 9.04, 0.48, LIGHT_BLUE, "B8CCE4")
                add_text(slide, items[-1], 0.68, 4.63, 8.64, 0.31, 14, NAVY, True, PP_ALIGN.CENTER)
            elif layout in {"image_right", "image_left"} and paths:
                if layout == "image_right":
                    add_bullets(slide, items, 0.48, 1.45, 3.35, 3.62, BLUE, 14.5)
                    add_image_fit(slide, paths[0], 4.05, 1.45, 5.47, 3.62)
                else:
                    add_image_fit(slide, paths[0], 0.48, 1.45, 5.47, 3.62)
                    add_bullets(slide, items, 6.17, 1.45, 3.35, 3.62, ORANGE, 14.5)
            elif layout == "three_images" and paths:
                gap=.18; w=(9.04-2*gap)/3
                for i,p in enumerate(paths[:3]): add_image_fit(slide,p,0.48+i*(w+gap),1.52,w,2.75)
                add_rect(slide,0.48,4.48,9.04,.56,LIGHT_BLUE,"B8CCE4")
                add_text(slide,items[-1],.68,4.58,8.64,.3,15,NAVY,True,PP_ALIGN.CENTER)
            elif layout == "cards": add_cards(slide, items)
            elif layout in {"process","timeline"}: add_process(slide, items)
            elif layout in {"two_columns","supported_limited","fact_hypothesis","tradeoff","status","matrix"}:
                mid = (len(items)+1)//2
                add_cards(slide, ["左侧｜"+"\n".join(items[:mid]), "右侧｜"+"\n".join(items[mid:])], h=3.48)
            elif layout == "kpi":
                add_cards(slide, items, h=3.45)
            elif layout == "trend":
                add_bullets(slide, items, 0.48, 1.48, 9.04, 3.55, BLUE, 15)
            elif layout == "diagram":
                add_process(slide, items)
            else:
                add_cards(slide, items)
            add_footer(slide, source)
        add_notes(slide, note)
    prs.core_properties.title = "Trench 与 SGT VDMOS 总剂量辐照效应对比及 TCAD 建模验证"
    prs.core_properties.subject = "答辩完整版"
    prs.core_properties.author = "项目研究团队"
    prs.save(PPTX)

    with (OUT / "素材来源清单.csv").open("w", newline="", encoding="utf-8-sig") as f:
        writer = csv.writer(f)
        writer.writerow(["页码","页面标题","相对路径","用途/布局","状态","存在","SHA-256"])
        writer.writerows(assets_log)

    mapping = {
        "final_pages": 74,
        "preserved_from_original": [1,9,10,11,12,13,14,15,16,17,27,28,29,30,31,32,33,35,36,37,41,42,44,45,46,47,48,50,51,67],
        "pages": [{"page":p[0],"title":p[1],"source":p[3],"layout":p[5],"notes":bool(p[6])} for p in PAGES]
    }
    (WORK / "page_manifest.json").write_text(json.dumps(mapping, ensure_ascii=False, indent=2), encoding="utf-8")
    print(PPTX)


if __name__ == "__main__":
    build()