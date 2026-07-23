# -*- coding: utf-8 -*-
import csv, json, re
from pathlib import Path
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation

OUT=Path(r"E:\仿真数据处理\答辩\PPT制作交付包\09_成品PPT")
PPTX=OUT/"S_T_VDMOS_TID_答辩完整版.pptx"
PREV=OUT/"逐页预览"

def num(p):
    m=re.search(r"(\d+)",p.stem)
    return int(m.group(1)) if m else 9999

imgs=sorted(PREV.glob("*.PNG"),key=num)
thumb_w,thumb_h=360,203
cols=5; rows=(len(imgs)+cols-1)//cols
sheet=Image.new("RGB",(cols*(thumb_w+18)+18,rows*(thumb_h+42)+18),"#E9EEF5")
d=ImageDraw.Draw(sheet)
for i,p in enumerate(imgs):
    im=Image.open(p).convert("RGB")
    im.thumbnail((thumb_w,thumb_h),Image.Resampling.LANCZOS)
    x=18+(i%cols)*(thumb_w+18); y=18+(i//cols)*(thumb_h+42)
    sheet.paste(im,(x,y))
    d.rectangle((x-1,y-1,x+thumb_w,y+thumb_h),outline="#AAB7C4",width=1)
    d.text((x,y+thumb_h+7),f"{i+1:02d}",fill="#243B53")
sheet.save(OUT/"contact_sheet.png",quality=92)

# 分批低码率接触表用于逐轮视觉检查，不作为正式交付替代品。
for batch_start in range(0, len(imgs), 15):
    batch = imgs[batch_start:batch_start+15]
    bw,bh,bcols = 190,107,3
    brows=(len(batch)+bcols-1)//bcols
    canvas=Image.new("RGB",(bcols*(bw+10)+10,brows*(bh+26)+10),"white")
    cd=ImageDraw.Draw(canvas)
    for j,p in enumerate(batch):
        im=Image.open(p).convert("RGB"); im.thumbnail((bw,bh),Image.Resampling.LANCZOS)
        xx=10+(j%bcols)*(bw+10); yy=10+(j//bcols)*(bh+26)
        canvas.paste(im,(xx,yy)); cd.text((xx,yy+bh+3),f"{batch_start+j+1:02d}",fill="#243B53")
    canvas.save(OUT/"制作工程"/f"contact_{batch_start//15+1:02d}.jpg",quality=18,optimize=True)

prs=Presentation(PPTX)
W,H=prs.slide_width,prs.slide_height
bounds=[]; tiny=[]; empty_notes=[]; empty_slides=[]
for si,s in enumerate(prs.slides,1):
    nonempty=0
    if not s.has_notes_slide or not s.notes_slide.notes_text_frame.text.strip(): empty_notes.append(si)
    for sh in s.shapes:
        if sh.left < 0 or sh.top < 0 or sh.left+sh.width > W+1000 or sh.top+sh.height > H+1000:
            bounds.append({"slide":si,"shape":sh.name})
        if getattr(sh,"has_text_frame",False) and sh.text.strip():
            nonempty += 1
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    if r.text.strip() and r.font.size and r.font.size.pt < 8:
                        tiny.append({"slide":si,"shape":sh.name,"size":r.font.size.pt,"text":r.text[:40]})
    if nonempty==0: empty_slides.append(si)

missing=[]
with (OUT/"素材来源清单.csv").open(encoding="utf-8-sig") as f:
    for row in csv.DictReader(f):
        if row["存在"].lower() not in {"true","1"}: missing.append(row)

metrics={
 "slide_count":len(prs.slides),"preview_count":len(imgs),"notes_count":len(prs.slides)-len(empty_notes),
 "out_of_bounds":bounds,"font_under_8pt":tiny,"empty_notes":empty_notes,"empty_text_slides":empty_slides,
 "missing_assets":missing,"contact_sheet_size":sheet.size,
 "pptx_bytes":PPTX.stat().st_size,"pdf_exists":(OUT/"S_T_VDMOS_TID_答辩完整版.pdf").exists()
}
(OUT/"制作工程"/"qa_metrics.json").write_text(json.dumps(metrics,ensure_ascii=False,indent=2),encoding="utf-8")
print(json.dumps({k:(len(v) if isinstance(v,list) else v) for k,v in metrics.items()},ensure_ascii=False))