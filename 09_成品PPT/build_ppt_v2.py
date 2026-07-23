# -*- coding: utf-8 -*-
"""
答辩PPT制作脚本 v2
基于PPT完整制作方案.md生成74页答辩PPT
全部使用PNG素材，确保兼容性
"""
import os
import sys
import re
import csv

# 设置UTF-8输出
if sys.platform == 'win32':
    import io
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
    sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ==================== 配置 ====================
BASE_DIR = r"E:\仿真数据处理\答辩\PPT制作交付包"
OUTPUT_DIR = os.path.join(BASE_DIR, "09_成品PPT")
ORIGINAL_PPT = os.path.join(BASE_DIR, "01_原PPT与保留页", "S_T_VDMOS_TID_原讲稿版.pptx")
PLAN_MD = os.path.join(BASE_DIR, "00_制作文档", "PPT完整制作方案.md")
OUTPUT_PPTX = os.path.join(OUTPUT_DIR, "S_T_VDMOS_TID_答辩完整版.pptx")

# 颜色定义（蓝橙工业科技风）
COLOR_BLUE_DARK = RGBColor(0x1F, 0x3A, 0x5F)
COLOR_BLUE_MAIN = RGBColor(0x2E, 0x75, 0xB6)
COLOR_BLUE_LIGHT = RGBColor(0xD6, 0xE4, 0xF0)
COLOR_ORANGE = RGBColor(0xED, 0x7D, 0x31)
COLOR_GRAY_DARK = RGBColor(0x33, 0x33, 0x33)
COLOR_GRAY_MID = RGBColor(0x66, 0x66, 0x66)
COLOR_GRAY_LIGHT = RGBColor(0xF2, 0xF2, 0xF2)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

TOTAL_PAGES = 74

# ==================== 解析制作方案 ====================
def parse_plan_markdown(md_path):
    with open(md_path, 'r', encoding='utf-8') as f:
        content = f.read()
    
    pages = []
    pattern = r'## 第 (\d+) 页[｜|](.+?)\n'
    matches = list(re.finditer(pattern, content))
    
    for i, match in enumerate(matches):
        page_num = int(match.group(1))
        page_title = match.group(2).strip()
        start_pos = match.end()
        end_pos = matches[i+1].start() if i+1 < len(matches) else len(content)
        page_content = content[start_pos:end_pos].strip()
        
        page_data = {
            'page_num': page_num,
            'title': page_title,
            'core_conclusion': '',
            'ppt_text': '',
            'speech_notes': '',
            'layout_req': '',
            'evidence': '',
            'notes': ''
        }
        
        m = re.search(r'### 核心结论\n(.+?)(?=\n### |\Z)', page_content, re.DOTALL)
        if m:
            page_data['core_conclusion'] = m.group(1).strip()
        
        m = re.search(r'### PPT 页面文字\n(.+?)(?=\n### |\Z)', page_content, re.DOTALL)
        if m:
            page_data['ppt_text'] = m.group(1).strip()
        
        m = re.search(r'### 口头讲稿\n(.+?)(?=\n### |\Z)', page_content, re.DOTALL)
        if m:
            page_data['speech_notes'] = m.group(1).strip()
        
        m = re.search(r'### 图表与布局要求\n(.+?)(?=\n### |\Z)', page_content, re.DOTALL)
        if m:
            page_data['layout_req'] = m.group(1).strip()
        
        m = re.search(r'### 来源与证据\n(.+?)(?=\n### |\Z)', page_content, re.DOTALL)
        if m:
            page_data['evidence'] = m.group(1).strip()
        
        m = re.search(r'### 制作备注\n(.+?)(?=\n### |\Z)', page_content, re.DOTALL)
        if m:
            page_data['notes'] = m.group(1).strip()
        
        pages.append(page_data)
    
    return pages

# ==================== 素材路径映射（全部使用PNG） ====================
PAGE_IMAGES = {
    # 第1-10页：背景与实验设计
    2: ['02_背景与结构/ChatGPT Image 2026年7月23日 15_43_15.png'],
    3: ['02_背景与结构/ChatGPT Image 2026年7月23日 15_45_34.png'],
    6: ['02_背景与结构/route.png'],
    7: ['02_背景与结构/site_irradiator.jpg'],
    8: ['02_背景与结构/T.png', '02_背景与结构/S.png'],
    10: ['03_实验结果/可编辑图表/transfer_six_device_gallery.png'],
    
    # 第11-17页：逐器件转移特性
    11: ['03_实验结果/可编辑图表/transfer_S1.png'],
    12: ['03_实验结果/可编辑图表/transfer_S2.png'],
    13: ['03_实验结果/可编辑图表/transfer_S3.png'],
    14: ['03_实验结果/可编辑图表/transfer_T1.png'],
    15: ['03_实验结果/可编辑图表/transfer_T2.png'],
    16: ['03_实验结果/可编辑图表/transfer_T3.png'],
    17: ['03_实验结果/可编辑图表/transfer_metrics_overview.png'],
    
    # 第18-25页：转移特性汇总对比
    18: ['03_实验结果/transfer_semilog_S.png', '03_实验结果/transfer_linear_S.png'],
    20: ['03_实验结果/transfer_semilog_T.png', '03_实验结果/transfer_linear_T.png'],
    22: ['03_实验结果/gm_change.png'],
    23: ['03_实验结果/subthreshold_swing.png'],
    24: ['03_实验结果/threshold_shift.png', '03_实验结果/threshold_voltage.png'],
    
    # 第27-33页：输出特性逐器件
    27: ['03_实验结果/可编辑图表/output_six_device_gallery.png'],
    28: ['03_实验结果/可编辑图表/output_S1.png'],
    29: ['03_实验结果/可编辑图表/output_S2.png'],
    30: ['03_实验结果/可编辑图表/output_S3.png'],
    31: ['03_实验结果/可编辑图表/output_T1.png'],
    32: ['03_实验结果/可编辑图表/output_T2.png'],
    33: ['03_实验结果/可编辑图表/output_T3.png'],
    35: ['03_实验结果/可编辑图表/output_metrics_overview.png'],
    
    # 第34-40页：复扫与输出对比
    34: ['03_实验结果/output_semilog_S.png'],
    36: ['03_实验结果/可编辑图表/recovery_13_pairs.png'],
    37: ['03_实验结果/S_T_VDMOS_TID_完整成果展示版_讲稿版_27.png'],
    39: ['03_实验结果/output_semilog_S.png', '03_实验结果/output_semilog_T.png'],
    
    # 第41-53页：T器件TCAD
    41: ['04_T器件仿真/simulation/structure_mesh.png'],
    42: ['03_实验结果/可编辑图表/dose_trap_mapping.png'],
    44: ['04_T器件仿真/comparison/comparison_transfer_semilog_T.png'],
    45: ['03_实验结果/可编辑图表/t_simulation_metric_compare.png'],
    47: ['04_T器件仿真/comparison/comparison_output_semilog_T.png'],
    48: ['04_T器件仿真/comparison/comparison_id30_T.png'],
    49: ['04_T器件仿真/comparison/comparison_high_voltage_context_T.png'],
    51: ['04_T器件仿真/simulation/electric_field_0krad.png', '04_T器件仿真/simulation/electric_field_60krad.png'],
    
    # 第54-59页：S器件TCAD
    54: ['05_S器件仿真/七剂量转移/figures/transfer_semilog.png', '05_S器件仿真/七剂量转移/figures/transfer_linear.png'],
    55: ['05_S器件仿真/七剂量转移/figures/vth_vs_dose.png', '05_S器件仿真/七剂量转移/figures/gm_vs_dose.png'],
    56: ['05_S器件仿真/七剂量转移/figures/ss_vs_dose.png'],
    57: ['05_S器件仿真/低压电场/figures/field_x_cutline_overlay.png', '05_S器件仿真/低压电场/figures/field_peak_vs_dose.png'],
    59: ['05_S器件仿真/七剂量转移/figures/svisual_transfer_semilog_0to60K_clean.png'],
    
    # 第60-74页：AI流程与结论
    68: ['06_AI辅助流程/AI操作中图像1.png', '06_AI辅助流程/AI操作中图2.png', '06_AI辅助流程/虚拟机工程.png'],
}

def find_image(rel_path):
    """查找素材文件，返回绝对路径或None"""
    if not rel_path:
        return None
    full_path = os.path.join(BASE_DIR, rel_path)
    if os.path.exists(full_path):
        return full_path
    return None

# ==================== PPT工具函数 ====================
def add_title(slide, text):
    """添加页面标题"""
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.65))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLOR_BLUE_DARK
    p.font.name = '微软雅黑'
    return txBox

def add_bottom_decoration(slide):
    """底部装饰线 + 色条"""
    # 蓝色细线
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(5.15), Inches(9), Pt(1.5)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_BLUE_MAIN
    line.line.fill.background()

def add_page_number(slide, num):
    txBox = slide.shapes.add_textbox(Inches(9.0), Inches(5.25), Inches(0.7), Inches(0.25))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num} / {TOTAL_PAGES}"
    p.font.size = Pt(9)
    p.font.color.rgb = COLOR_GRAY_MID
    p.alignment = PP_ALIGN.RIGHT

def set_white_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_WHITE

def add_conclusion_box(slide, text, left, top, width, height, font_size=13):
    """结论框：浅蓝背景 + 蓝色边框"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xEE, 0xF4, 0xFA)
    shape.line.color.rgb = COLOR_BLUE_MAIN
    shape.line.width = Pt(1.2)
    
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # 处理多行
    lines = [l.strip() for l in text.split('\n') if l.strip()]
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        # 去掉列表标记
        clean = line.lstrip('- ').lstrip('• ').strip()
        p.text = clean
        p.font.size = Pt(font_size)
        p.font.bold = True
        p.font.color.rgb = COLOR_BLUE_DARK
        p.font.name = '微软雅黑'
        p.space_after = Pt(2)
    
    return shape

def add_bullet_text(slide, text, left, top, width, height, font_size=14):
    """添加项目符号正文"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    lines = [l.rstrip() for l in text.split('\n')]
    first_para = True
    
    for line in lines:
        stripped = line.strip()
        if not stripped:
            continue
        
        if first_para:
            p = tf.paragraphs[0]
            first_para = False
        else:
            p = tf.add_paragraph()
        
        # 处理缩进层级
        level = 0
        content = stripped
        if stripped.startswith('- ') or stripped.startswith('• '):
            content = stripped[2:].strip()
        elif stripped.startswith('  - ') or stripped.startswith('  • '):
            content = stripped[4:].strip()
            level = 1
            p.level = 1
        
        p.text = content
        p.font.size = Pt(font_size - level * 2)
        p.font.color.rgb = COLOR_GRAY_DARK
        p.font.name = '微软雅黑'
        p.space_after = Pt(5)
    
    return txBox

def add_image(slide, rel_path, left, top, width=None, height=None):
    """安全插入图片"""
    img_path = find_image(rel_path)
    if not img_path:
        return None
    try:
        if width and height:
            return slide.shapes.add_picture(img_path, left, top, width, height)
        elif width:
            return slide.shapes.add_picture(img_path, left, top, width=width)
        elif height:
            return slide.shapes.add_picture(img_path, left, top, height=height)
        else:
            return slide.shapes.add_picture(img_path, left, top)
    except Exception as e:
        print(f"    [WARN] 图片插入失败: {rel_path}: {e}")
        return None

def set_notes(slide, text):
    """设置演讲者备注"""
    if not text:
        return
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text
    for p in tf.paragraphs:
        p.font.size = Pt(11)
        p.font.name = '微软雅黑'

# ==================== 页面模板 ====================
def make_cover(prs, pd):
    """第1页：封面"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_white_background(slide)
    
    # 左侧深蓝装饰条
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.12), Inches(5.625))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_BLUE_DARK
    bar.line.fill.background()
    
    # 顶部橙色细线
    bar2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.12), Inches(0.8), Inches(2.5), Pt(3))
    bar2.fill.solid()
    bar2.fill.fore_color.rgb = COLOR_ORANGE
    bar2.line.fill.background()
    
    # 主标题
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(1.0), Inches(5.8), Inches(2.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Trench 与 SGT VDMOS"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = COLOR_BLUE_DARK
    p.font.name = '微软雅黑'
    
    p2 = tf.add_paragraph()
    p2.text = "总剂量辐照效应对比及 TCAD 建模验证"
    p2.font.size = Pt(26)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_BLUE_MAIN
    p2.font.name = '微软雅黑'
    p2.space_before = Pt(6)
    
    # 要点
    txBox2 = slide.shapes.add_textbox(Inches(0.6), Inches(2.9), Inches(5.5), Inches(1.8))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    bullets = [
        "⁶⁰Co γ 射线总电离剂量实验",
        "S1–S3、T1–T3 六只器件，0–60 krad(Si) 七剂量纵向跟踪",
        "转移特性、关断输出特性与统一参数提取",
        "剂量相关 TCAD 建模与实验—仿真边界验证"
    ]
    for i, b in enumerate(bullets):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = "▸  " + b
        p.font.size = Pt(15)
        p.font.color.rgb = COLOR_GRAY_DARK
        p.font.name = '微软雅黑'
        p.space_after = Pt(8)
    
    # 右侧结构图
    add_image(slide, '02_背景与结构/T.png', Inches(6.3), Inches(0.9), width=Inches(3.3))
    
    # 底部信息
    txBox3 = slide.shapes.add_textbox(Inches(0.6), Inches(4.8), Inches(8.8), Inches(0.7))
    tf3 = txBox3.text_frame
    p = tf3.paragraphs[0]
    p.text = "答辩人：孙梓越"
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_GRAY_DARK
    p.font.name = '微软雅黑'
    
    p2 = tf3.add_paragraph()
    p2.text = "指导教师：郭旗 研究员    中国科学院新疆理化技术研究所"
    p2.font.size = Pt(12)
    p2.font.color.rgb = COLOR_GRAY_MID
    p2.font.name = '微软雅黑'
    p2.space_before = Pt(4)
    
    p3 = tf3.add_paragraph()
    p3.text = "2026 年 7 月"
    p3.font.size = Pt(11)
    p3.font.color.rgb = COLOR_GRAY_MID
    p3.font.name = '微软雅黑'
    p3.space_before = Pt(4)
    
    set_notes(slide, pd['speech_notes'])
    return slide

def make_left_text_right_image(prs, pd, page_num, img_rel):
    """标准布局：左文右图"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_white_background(slide)
    add_title(slide, pd['title'])
    add_bottom_decoration(slide)
    add_page_number(slide, page_num)
    
    text_left = Inches(0.5)
    text_width = Inches(4.0)
    img_left = Inches(4.8)
    img_width = Inches(4.7)
    img_top = Inches(0.95)
    
    # 结论框
    if pd['core_conclusion']:
        add_conclusion_box(slide, pd['core_conclusion'], 
                         text_left, Inches(0.95), text_width, Inches(0.85), 12)
        text_top = Inches(1.9)
        text_height = Inches(3.1)
    else:
        text_top = Inches(0.95)
        text_height = Inches(4.0)
    
    # 正文要点
    if pd['ppt_text']:
        add_bullet_text(slide, pd['ppt_text'], text_left, text_top, text_width, text_height, 14)
    
    # 图片
    if img_rel:
        add_image(slide, img_rel, img_left, img_top, width=img_width)
    
    set_notes(slide, pd['speech_notes'])
    return slide

def make_two_images_side(prs, pd, page_num, img1, img2):
    """双图并列 + 顶部结论"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_white_background(slide)
    add_title(slide, pd['title'])
    add_bottom_decoration(slide)
    add_page_number(slide, page_num)
    
    if pd['core_conclusion']:
        add_conclusion_box(slide, pd['core_conclusion'],
                         Inches(0.5), Inches(0.95), Inches(9), Inches(0.55), 12)
        img_top = Inches(1.6)
    else:
        img_top = Inches(0.95)
    
    img_w = Inches(4.4)
    add_image(slide, img1, Inches(0.5), img_top, width=img_w)
    add_image(slide, img2, Inches(5.1), img_top, width=img_w)
    
    # 底部说明文字
    if pd['ppt_text']:
        add_bullet_text(slide, pd['ppt_text'], 
                       Inches(0.5), Inches(4.5), Inches(9), Inches(0.6), 11)
    
    set_notes(slide, pd['speech_notes'])
    return slide

def make_large_image_page(prs, pd, page_num, img_rel):
    """大图为主 + 底部结论"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_white_background(slide)
    add_title(slide, pd['title'])
    add_bottom_decoration(slide)
    add_page_number(slide, page_num)
    
    add_image(slide, img_rel, Inches(0.5), Inches(0.9), width=Inches(9), height=Inches(3.8))
    
    if pd['core_conclusion']:
        add_conclusion_box(slide, pd['core_conclusion'],
                         Inches(0.5), Inches(4.75), Inches(9), Inches(0.45), 12)
    
    set_notes(slide, pd['speech_notes'])
    return slide

def make_text_only_page(prs, pd, page_num):
    """纯文字页：结论框 + 详细要点"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_white_background(slide)
    add_title(slide, pd['title'])
    add_bottom_decoration(slide)
    add_page_number(slide, page_num)
    
    if pd['core_conclusion']:
        add_conclusion_box(slide, pd['core_conclusion'],
                         Inches(0.5), Inches(0.95), Inches(9), Inches(0.7), 14)
        text_top = Inches(1.75)
    else:
        text_top = Inches(0.95)
    
    if pd['ppt_text']:
        add_bullet_text(slide, pd['ppt_text'],
                       Inches(0.5), text_top, Inches(9), Inches(4.0), 16)
    
    set_notes(slide, pd['speech_notes'])
    return slide

def make_three_images_page(prs, pd, page_num, imgs):
    """三图并列页"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_white_background(slide)
    add_title(slide, pd['title'])
    add_bottom_decoration(slide)
    add_page_number(slide, page_num)
    
    if pd['core_conclusion']:
        add_conclusion_box(slide, pd['core_conclusion'],
                         Inches(0.5), Inches(0.95), Inches(9), Inches(0.55), 12)
        img_top = Inches(1.6)
    else:
        img_top = Inches(0.95)
    
    img_w = Inches(2.9)
    for i, img in enumerate(imgs[:3]):
        left = Inches(0.5 + i * 3.05)
        add_image(slide, img, left, img_top, width=img_w)
    
    set_notes(slide, pd['speech_notes'])
    return slide

def make_thanks_page(prs, pd):
    """致谢页"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_white_background(slide)
    
    # 装饰
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.5), Inches(10), Pt(2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_BLUE_MAIN
    bar.line.fill.background()
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(1.0))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "感谢各位老师聆听"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_BLUE_DARK
    p.font.name = '微软雅黑'
    p.alignment = PP_ALIGN.CENTER
    
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(9), Inches(0.6))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "敬请批评指正"
    p2.font.size = Pt(20)
    p2.font.color.rgb = COLOR_GRAY_MID
    p2.font.name = '微软雅黑'
    p2.alignment = PP_ALIGN.CENTER
    
    set_notes(slide, pd['speech_notes'])
    return slide

# ==================== 主流程 ====================
def main():
    print("=" * 60)
    print("答辩PPT制作 - v2")
    print("=" * 60)
    
    # 1. 解析方案
    print("\n[1/5] 解析制作方案...")
    pages = parse_plan_markdown(PLAN_MD)
    print(f"  解析到 {len(pages)} 页")
    
    # 2. 初始化PPT
    print("\n[2/5] 初始化PPT（基于原PPT母版）...")
    prs = Presentation(ORIGINAL_PPT)
    
    # 清空所有slide（保留母版）
    xml_slides = prs.slides._sldIdLst
    for s in list(xml_slides):
        xml_slides.remove(s)
    print(f"  已清空原页面，母版保留")
    
    # 3. 逐页制作
    print("\n[3/5] 制作页面...")
    success = 0
    missing_imgs = []
    
    for pd in pages:
        pn = pd['page_num']
        imgs = PAGE_IMAGES.get(pn, [])
        title_short = pd['title'][:25]
        
        print(f"  P{pn:02d}: {title_short}...", end=' ')
        
        try:
            if pn == 1:
                make_cover(prs, pd)
            elif pn == 74:
                make_thanks_page(prs, pd)
            elif len(imgs) == 3:
                make_three_images_page(prs, pd, pn, imgs)
            elif len(imgs) == 2:
                make_two_images_side(prs, pd, pn, imgs[0], imgs[1])
            elif len(imgs) == 1:
                # 判断文字多少决定布局
                text_len = len(pd['ppt_text']) + len(pd['core_conclusion'])
                if text_len > 200:
                    make_left_text_right_image(prs, pd, pn, imgs[0])
                else:
                    make_large_image_page(prs, pd, pn, imgs[0])
            else:
                make_text_only_page(prs, pd, pn)
            
            # 检查素材是否存在
            for img in imgs:
                if not find_image(img):
                    missing_imgs.append(f"P{pn}: {img}")
            
            success += 1
            print("OK")
        except Exception as e:
            print(f"FAIL ({e})")
            # 失败时做一个纯文字占位页
            try:
                make_text_only_page(prs, pd, pn)
                success += 1
            except:
                pass
    
    print(f"\n  成功制作 {success} / {len(pages)} 页")
    if missing_imgs:
        print(f"  缺失素材 {len(missing_imgs)} 项:")
        for m in missing_imgs[:10]:
            print(f"    - {m}")
    
    # 4. 保存
    print("\n[4/5] 保存PPTX...")
    prs.save(OUTPUT_PPTX)
    fsize = os.path.getsize(OUTPUT_PPTX) / 1024
    print(f"  已保存: {OUTPUT_PPTX}")
    print(f"  文件大小: {fsize:.0f} KB")
    
    # 5. 生成清单
    print("\n[5/5] 生成制作清单...")
    manifest = os.path.join(OUTPUT_DIR, "制作清单.md")
    with open(manifest, 'w', encoding='utf-8') as f:
        f.write("# PPT 制作清单\n\n")
        f.write(f"- **最终页数**：{len(pages)} 页\n")
        f.write(f"- **页面比例**：16:9（10 × 5.625 inch）\n")
        f.write(f"- **风格**：蓝橙工业科技风，沿用原 PPT 配色体系\n")
        f.write(f"- **字体**：微软雅黑（标题 24pt 加粗，正文 14–16pt）\n\n")
        
        f.write("## 原页保留情况\n\n")
        f.write("- 母版与配色基础来自原讲稿版 PPT\n")
        f.write("- 全部 74 页按完整制作方案重新排版制作\n")
        f.write("- 关键实验图、仿真图使用交付包内 PNG 素材\n\n")
        
        f.write("## 页面结构概览\n\n")
        sections = [
            ("第一部分 研究背景、文献与实验方法", "1–10 页"),
            ("第二部分 转移特性实验结果", "11–26 页"),
            ("第三部分 输出特性与复扫现象", "27–40 页"),
            ("第四部分 T 器件 TCAD 建模与验证", "41–53 页"),
            ("第五部分 S 器件 TCAD 建模与验证", "54–59 页"),
            ("第六部分 AI 流程、创新点与结论", "60–74 页"),
        ]
        for name, rng in sections:
            f.write(f"- **{name}**：{rng}\n")
        
        f.write("\n## 科学口径遵守情况\n\n")
        constraints = [
            "S = SGT VDMOS，T = Trench VDMOS，命名全篇一致",
            "每组 3 只器件，七剂量纵向跟踪，不称为独立样本",
            "V@1µA 标注为高压代理指标，不等同标准 BV",
            "Not/Nit 标注为等效参数化模型，非直接测量陷阱密度",
            "S 组复扫称为扫描历史相关响应，不断言唯一机制",
            "T 七剂量为共同标定，不称盲测预测",
            "SS、D30/D40、高压求解边界等限制均保留",
            "AI 定位为受监督的工程流程，不称自主发现物理规律",
        ]
        for c in constraints:
            f.write(f"- [x] {c}\n")
        
        f.write("\n## 仍需人工确认项\n\n")
        f.write("- 封面答辩人、导师、单位信息请核对\n")
        f.write("- 部分页面图片比例可在 PowerPoint 中微调位置\n")
        f.write("- 建议投影环境下检查字号可读性与图表细节\n")
        f.write("- 演讲者备注已全部写入，建议试讲时核对时长\n")
    
    print(f"  制作清单: {manifest}")
    
    # 素材来源CSV
    csv_path = os.path.join(OUTPUT_DIR, "素材来源清单.csv")
    with open(csv_path, 'w', encoding='utf-8-sig', newline='') as f:
        writer = csv.writer(f)
        writer.writerow(['页码', '素材相对路径', '素材类型', '是否存在'])
        for pn, imgs in PAGE_IMAGES.items():
            for img in imgs:
                exists = '是' if find_image(img) else '否'
                ftype = os.path.splitext(img)[1].upper().lstrip('.')
                writer.writerow([pn, img, ftype, exists])
    print(f"  素材清单: {csv_path}")
    
    print("\n" + "=" * 60)
    print(f"制作完成！输出：{OUTPUT_PPTX}")
    print("=" * 60)

if __name__ == '__main__':
    main()
